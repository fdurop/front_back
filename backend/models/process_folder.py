from neo4j import GraphDatabase
import pandas as pd
from typing import Optional, Dict, List, Tuple
import json
from collections import defaultdict


class Neo4jKnowledgeGraph:
    """Neo4jçŸ¥è¯†å›¾è°±è¿æ¥å™¨"""

    def __init__(self, uri: str, user: str, password: str):
        """åˆå§‹åŒ–Neo4jè¿æ¥"""
        self.driver = None
        try:
            config = {
                "keep_alive": True,
                "max_connection_lifetime": 3600,
                "max_connection_pool_size": 100
            }
            self.driver = GraphDatabase.driver(uri, auth=(user, password), **config)

            # æµ‹è¯•è¿æ¥
            with self.driver.session() as session:
                session.run("RETURN 1")
            print("âœ… Neo4jè¿æ¥æˆåŠŸ")

        except Exception as e:
            print(f"âŒ Neo4jè¿æ¥å¤±è´¥: {e}")
            raise

    def close(self):
        """å…³é—­è¿æ¥"""
        if self.driver:
            self.driver.close()

    def clear_database(self):
        """æ¸…ç©ºæ•°æ®åº“ï¼ˆè°¨æ…ä½¿ç”¨ï¼‰"""
        with self.driver.session() as session:
            session.run("MATCH (n) DETACH DELETE n")
        print("ğŸ—‘ï¸ æ•°æ®åº“å·²æ¸…ç©º")

    def create_document_node(self, doc_name: str, doc_type: str = "ppt", metadata: Dict = None):
        """åˆ›å»ºæ–‡æ¡£èŠ‚ç‚¹"""
        with self.driver.session() as session:
            query = """
            MERGE (d:Document {name: $doc_name})
            SET d.type = $doc_type
            SET d.created_at = datetime()
            """

            if metadata:
                for key, value in metadata.items():
                    query += f"SET d.{key} = ${key} "

            params = {"doc_name": doc_name, "doc_type": doc_type}
            if metadata:
                params.update(metadata)

            session.run(query, params)
        print(f"ğŸ“„ åˆ›å»ºæ–‡æ¡£èŠ‚ç‚¹: {doc_name}")

    def create_entity_node(self, entity_name: str, entity_type: str, description: str = "",
                           metadata: Dict = None):
        """åˆ›å»ºå®ä½“èŠ‚ç‚¹"""
        with self.driver.session() as session:
            query = """
            MERGE (e:Entity {name: $entity_name})
            SET e.type = $entity_type
            SET e.description = $description
            SET e.updated_at = datetime()
            """

            params = {
                "entity_name": entity_name,
                "entity_type": entity_type,
                "description": description
            }

            if metadata:
                for key, value in metadata.items():
                    if isinstance(value, (str, int, float, bool)):
                        query += f"SET e.{key} = ${key} "
                        params[key] = value

            session.run(query, params)

    def create_relationship(self, source_name: str, target_name: str, relation_type: str,
                            properties: Dict = None):
        """åˆ›å»ºå…³ç³»"""
        with self.driver.session() as session:
            # ç¡®ä¿æºèŠ‚ç‚¹å’Œç›®æ ‡èŠ‚ç‚¹å­˜åœ¨
            session.run("""
                MERGE (s:Entity {name: $source_name})
                MERGE (t:Entity {name: $target_name})
            """, source_name=source_name, target_name=target_name)

            # åˆ›å»ºå…³ç³»
            query = f"""
            MATCH (s:Entity {{name: $source_name}})
            MATCH (t:Entity {{name: $target_name}})
            MERGE (s)-[r:{relation_type}]->(t)
            SET r.created_at = datetime()
            """

            params = {"source_name": source_name, "target_name": target_name}

            if properties:
                for key, value in properties.items():
                    if isinstance(value, (str, int, float, bool)):
                        query += f"SET r.{key} = ${key} "
                        params[key] = value

            session.run(query, params)

    def save_extracted_data(self, extracted_data, ppt_name: str):
        """ä¿å­˜æŠ½å–çš„å®ä½“å…³ç³»æ•°æ®åˆ°Neo4j"""
        try:
            print(f"ğŸ’¾ å¼€å§‹ä¿å­˜æ•°æ®åˆ°Neo4j: {ppt_name}")

            # 1. åˆ›å»ºPPTæ–‡æ¡£èŠ‚ç‚¹
            self.create_document_node(ppt_name, "ppt", {
                "total_entities": len(extracted_data.entities),
                "total_relationships": len(extracted_data.relationships)
            })

            # 2. æ‰¹é‡åˆ›å»ºå®ä½“èŠ‚ç‚¹
            print(f"   åˆ›å»º {len(extracted_data.entities)} ä¸ªå®ä½“èŠ‚ç‚¹...")
            entity_count = 0
            for entity in extracted_data.entities:
                try:
                    metadata = {k: v for k, v in entity.items()
                                if k not in ['name', 'type', 'description'] and
                                isinstance(v, (str, int, float, bool))}

                    self.create_entity_node(
                        entity['name'],
                        entity.get('type', 'unknown'),
                        entity.get('description', ''),
                        metadata
                    )

                    # åˆ›å»ºæ–‡æ¡£åˆ°å®ä½“çš„åŒ…å«å…³ç³»
                    self.create_relationship(ppt_name, entity['name'], "CONTAINS")
                    entity_count += 1

                except Exception as e:
                    print(f"     âš ï¸ åˆ›å»ºå®ä½“å¤±è´¥ {entity['name']}: {e}")
                    continue

            # 3. æ‰¹é‡åˆ›å»ºå…³ç³»
            print(f"   åˆ›å»º {len(extracted_data.relationships)} ä¸ªå…³ç³»...")
            relation_count = 0
            for rel in extracted_data.relationships:
                try:
                    properties = {k: v for k, v in rel.items()
                                  if k not in ['source', 'target', 'relation'] and
                                  isinstance(v, (str, int, float, bool))}

                    self.create_relationship(
                        rel['source'],
                        rel['target'],
                        rel['relation'].upper().replace(' ', '_'),
                        properties
                    )
                    relation_count += 1

                except Exception as e:
                    print(f"     âš ï¸ åˆ›å»ºå…³ç³»å¤±è´¥ {rel['source']}->{rel['target']}: {e}")
                    continue

            print(f"âœ… æ•°æ®ä¿å­˜å®Œæˆ:")
            print(f"   ğŸ“„ æ–‡æ¡£: {ppt_name}")
            print(f"   ğŸ·ï¸  å®ä½“: {entity_count}/{len(extracted_data.entities)}")
            print(f"   ğŸ”— å…³ç³»: {relation_count}/{len(extracted_data.relationships)}")

            return {
                'document': ppt_name,
                'entities_saved': entity_count,
                'relationships_saved': relation_count,
                'success': True
            }

        except Exception as e:
            print(f"âŒ ä¿å­˜æ•°æ®å¤±è´¥: {e}")
            return {
                'document': ppt_name,
                'entities_saved': 0,
                'relationships_saved': 0,
                'success': False,
                'error': str(e)
            }

    def query_entities(self, limit: int = 10) -> List[Dict]:
        """æŸ¥è¯¢å®ä½“"""
        with self.driver.session() as session:
            result = session.run("""
                MATCH (e:Entity)
                RETURN e.name as name, e.type as type, e.description as description
                LIMIT $limit
            """, limit=limit)

            return [dict(record) for record in result]

    def query_relationships(self, limit: int = 10) -> List[Dict]:
        """æŸ¥è¯¢å…³ç³»"""
        with self.driver.session() as session:
            result = session.run("""
                MATCH (s:Entity)-[r]->(t:Entity)
                RETURN s.name as source, type(r) as relation, t.name as target
                LIMIT $limit
            """, limit=limit)

            return [dict(record) for record in result]

    def get_statistics(self) -> Dict:
        """è·å–æ•°æ®åº“ç»Ÿè®¡ä¿¡æ¯"""
        with self.driver.session() as session:
            # èŠ‚ç‚¹ç»Ÿè®¡
            node_result = session.run("MATCH (n) RETURN count(n) as total_nodes")
            total_nodes = node_result.single()['total_nodes']

            # å…³ç³»ç»Ÿè®¡
            rel_result = session.run("MATCH ()-[r]->() RETURN count(r) as total_relationships")
            total_relationships = rel_result.single()['total_relationships']

            # å®ä½“ç±»å‹ç»Ÿè®¡
            entity_type_result = session.run("""
                MATCH (e:Entity)
                RETURN e.type as entity_type, count(e) as count
                ORDER BY count DESC
            """)
            entity_types = [dict(record) for record in entity_type_result]

            # æ–‡æ¡£ç»Ÿè®¡
            doc_result = session.run("MATCH (d:Document) RETURN count(d) as total_documents")
            total_documents = doc_result.single()['total_documents']

            return {
                'total_nodes': total_nodes,
                'total_relationships': total_relationships,
                'total_documents': total_documents,
                'entity_types': entity_types
            }

    def search_entities_by_name(self, name_pattern: str, limit: int = 10) -> List[Dict]:
        """æŒ‰åç§°æœç´¢å®ä½“"""
        with self.driver.session() as session:
            result = session.run("""
                MATCH (e:Entity)
                WHERE toLower(e.name) CONTAINS toLower($pattern)
                RETURN e.name as name, e.type as type, e.description as description
                LIMIT $limit
            """, pattern=name_pattern, limit=limit)

            return [dict(record) for record in result]

    def get_entity_neighbors(self, entity_name: str, depth: int = 1) -> Dict:
        """è·å–å®ä½“çš„é‚»å±…èŠ‚ç‚¹"""
        with self.driver.session() as session:
            result = session.run(f"""
                MATCH path = (e:Entity {{name: $entity_name}})-[*1..{depth}]-(neighbor)
                RETURN neighbor.name as name, neighbor.type as type, 
                       neighbor.description as description
                LIMIT 20
            """, entity_name=entity_name)

            neighbors = [dict(record) for record in result]

            # è·å–ç›¸å…³å…³ç³»
            rel_result = session.run("""
                MATCH (e:Entity {name: $entity_name})-[r]-(neighbor)
                RETURN neighbor.name as neighbor, type(r) as relation,
                       CASE WHEN startNode(r).name = $entity_name 
                            THEN 'outgoing' ELSE 'incoming' END as direction
            """, entity_name=entity_name)

            relationships = [dict(record) for record in rel_result]

            return {
                'entity': entity_name,
                'neighbors': neighbors,
                'relationships': relationships
            }


def save_to_neo4j(extracted_data, ppt_name: str, neo4j_uri: str, neo4j_user: str, neo4j_password: str):
    """ä¿å­˜æ•°æ®åˆ°Neo4jçš„ä¸»å‡½æ•°"""
    kg = Neo4jKnowledgeGraph(neo4j_uri, neo4j_user, neo4j_password)

    try:
        result = kg.save_extracted_data(extracted_data, ppt_name)

        # æ˜¾ç¤ºç»Ÿè®¡ä¿¡æ¯
        stats = kg.get_statistics()
        print(f"\nğŸ“Š æ•°æ®åº“ç»Ÿè®¡ä¿¡æ¯:")
        print(f"   ğŸ“„ æ–‡æ¡£æ€»æ•°: {stats['total_documents']}")
        print(f"   ğŸ·ï¸  èŠ‚ç‚¹æ€»æ•°: {stats['total_nodes']}")
        print(f"   ğŸ”— å…³ç³»æ€»æ•°: {stats['total_relationships']}")

        # æ˜¾ç¤ºå®ä½“ç±»å‹åˆ†å¸ƒ
        print(f"\nğŸ·ï¸  å®ä½“ç±»å‹åˆ†å¸ƒ:")
        for entity_type in stats['entity_types'][:5]:
            print(f"   - {entity_type['entity_type']}: {entity_type['count']}ä¸ª")

        return result

    finally:
        kg.close()


import json
import os
import re
from typing import List, Dict, Tuple, Optional
from dataclasses import dataclass
from collections import defaultdict
import requests
import time


@dataclass
class ExtractedTriple:
    entities: List[Dict]
    relationships: List[Dict]
    attributes: List[Dict]


class DeepSeekClient:
    """DeepSeek APIå®¢æˆ·ç«¯"""

    def __init__(self, api_key, base_url="https://api.deepseek.com/v1", model="deepseek-chat"):
        self.api_key = api_key
        self.base_url = base_url
        self.model = model
        self.headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }

    def chat_completions_create(self, messages: List[Dict], temperature: float = 0.1, max_tokens: int = 1024) -> Dict:
        """è°ƒç”¨DeepSeek API"""
        try:
            response = requests.post(
                f"{self.base_url}/chat/completions",
                headers=self.headers,
                json={
                    "model": self.model,
                    "messages": messages,
                    "temperature": temperature,
                    "max_tokens": max_tokens
                },
                timeout=30
            )
            response.raise_for_status()
            return response.json()
        except Exception as e:
            print(f"DeepSeek APIè°ƒç”¨å¤±è´¥: {e}")
            return {"choices": [{"message": {"content": "{\"entities\": [], \"relationships\": []}"}}]}


class EntityExtractor:
    """å®ä½“å…³ç³»æŠ½å–å™¨"""

    def __init__(self, deepseek_api_key: str):
        self.deepseek = DeepSeekClient(deepseek_api_key)
        self.arduino_keywords = [
            'Arduino', 'LED', 'sensor', 'ä¼ æ„Ÿå™¨', 'pin', 'å¼•è„š', 'GPIO',
            'voltage', 'ç”µå‹', 'current', 'ç”µæµ', 'resistor', 'ç”µé˜»', 'PWM',
            'digital', 'æ•°å­—', 'analog', 'æ¨¡æ‹Ÿ', 'serial', 'ä¸²å£', 'I2C', 'SPI',
            'breadboard', 'é¢åŒ…æ¿', 'wire', 'å¯¼çº¿', 'ground', 'æ¥åœ°', 'VCC', '5V', '3.3V'
        ]

    def _extract_slide_number(self, filename: str) -> int:
        """ä»æ–‡ä»¶åä¸­æå–å¹»ç¯ç‰‡å·ç """
        match = re.search(r'slide_(\d+)', filename)
        return int(match.group(1)) if match else 0

    def load_multimodal_data(self, output_dir: str = "output") -> Dict:
        """åŠ è½½å¤šæ¨¡æ€é¢„å¤„ç†çš„è¾“å‡ºæ•°æ® - é€‚é…å®é™…æ–‡ä»¶æ ¼å¼"""
        result = {
            'slides': [],
            'images': []
        }

        try:
            # å®šä¹‰å­ç›®å½•è·¯å¾„
            text_dir = os.path.join(output_dir, "text")
            image_dir = os.path.join(output_dir, "images")

            # 1. åŠ è½½å¹»ç¯ç‰‡æ–‡æœ¬æ•°æ® (ä»textç›®å½•)
            if os.path.exists(text_dir):
                text_files = os.listdir(text_dir)
                slide_files = [f for f in text_files if
                               '_slide_' in f and f.endswith('.json') and not f.endswith('_desc.json')]

                for slide_file in slide_files:
                    slide_path = os.path.join(text_dir, slide_file)  # æ³¨æ„è¿™é‡Œæ”¹ä¸ºtext_dir
                    try:
                        with open(slide_path, 'r', encoding='utf-8') as f:
                            slide_data = json.load(f)

                        slide_num = self._extract_slide_number(slide_file)

                        result['slides'].append({
                            "slide_number": slide_num,
                            "content": slide_data,
                            "source_file": slide_file
                        })

                    except Exception as e:
                        print(f"âš ï¸ åŠ è½½å¹»ç¯ç‰‡æ–‡ä»¶å¤±è´¥ {slide_file}: {e}")

            # 2. åŠ è½½å›¾ç‰‡æ•°æ® (ä»imageç›®å½•)
            if os.path.exists(image_dir):
                image_files_list = os.listdir(image_dir)
                image_files = [f for f in image_files_list if f.endswith('.png') or f.endswith('.jpg')]

                for image_file in image_files:
                    # æŸ¥æ‰¾å¯¹åº”çš„æè¿°æ–‡ä»¶
                    desc_file = image_file.replace('.png', '_desc.json').replace('.jpg', '_desc.json')
                    desc_path = os.path.join(image_dir, desc_file)  # æ³¨æ„è¿™é‡Œæ”¹ä¸ºimage_dir

                    slide_num = self._extract_slide_number(image_file)

                    image_data = {
                        "image_path": os.path.join(image_dir, image_file),  # æ³¨æ„è¿™é‡Œæ”¹ä¸ºimage_dir
                        "slide_number": slide_num,
                        "filename": image_file,
                        "descriptions": [],
                        "ocr_text": ""
                    }

                    # å¦‚æœæœ‰æè¿°æ–‡ä»¶ï¼ŒåŠ è½½æè¿°ä¿¡æ¯
                    if os.path.exists(desc_path):
                        try:
                            with open(desc_path, 'r', encoding='utf-8') as f:
                                desc_data = json.load(f)
                                image_data["descriptions"] = desc_data.get("clip_descriptions", [])
                        except Exception as e:
                            print(f"âš ï¸ åŠ è½½å›¾ç‰‡æè¿°å¤±è´¥ {desc_file}: {e}")

                    result['images'].append(image_data)

            print(f"âœ… æ•°æ®åŠ è½½å®Œæˆ:")
            print(f"   - å¹»ç¯ç‰‡: {len(result['slides'])}ä¸ªæ–‡ä»¶")
            print(f"   - å›¾ç‰‡: {len(result['images'])}ä¸ª")

        except Exception as e:
            print(f"âŒ æ•°æ®åŠ è½½å¤±è´¥: {e}")

        return result

    def extract_entities_from_multimodal(self, multimodal_data: Dict) -> ExtractedTriple:
        """ä»å¤šæ¨¡æ€æ•°æ®ä¸­æŠ½å–å®ä½“å…³ç³»"""
        all_entities = []
        all_relationships = []
        all_attributes = []

        print("ğŸ” å¼€å§‹å®ä½“å…³ç³»æŠ½å–...")

        # 1. å¤„ç†å¹»ç¯ç‰‡æ–‡æœ¬å†…å®¹
        slides = multimodal_data.get('slides', [])
        for i, slide in enumerate(slides):
            print(f"   å¤„ç†å¹»ç¯ç‰‡ {i + 1}/{len(slides)}: {slide.get('source_file', '')}")
            slide_entities, slide_relations = self._extract_from_slide_text(slide)
            all_entities.extend(slide_entities)
            all_relationships.extend(slide_relations)
            time.sleep(0.5)  # é¿å…APIè°ƒç”¨è¿‡å¿«

        # 2. å¤„ç†å›¾ç‰‡å†…å®¹
        images = multimodal_data.get('images', [])
        for i, image_data in enumerate(images):
            print(f"   å¤„ç†å›¾ç‰‡ {i + 1}/{len(images)}: {image_data.get('filename', '')}")
            img_entities = self._extract_from_image(image_data)
            all_entities.extend(img_entities)

        # å»é‡å¤„ç†
        all_entities = self._deduplicate_entities(all_entities)
        all_relationships = self._deduplicate_relationships(all_relationships)

        print(f"âœ… å®ä½“å…³ç³»æŠ½å–å®Œæˆ: {len(all_entities)}ä¸ªå®ä½“, {len(all_relationships)}ä¸ªå…³ç³»")

        return ExtractedTriple(
            entities=all_entities,
            relationships=all_relationships,
            attributes=all_attributes
        )

    def _extract_from_slide_text(self, slide: Dict) -> Tuple[List[Dict], List[Dict]]:
        """ä»å¹»ç¯ç‰‡æ–‡æœ¬ä¸­æŠ½å–å®ä½“å’Œå…³ç³»"""
        slide_content = slide.get('content', {})
        slide_num = slide.get('slide_number', 0)

        # æå–æ–‡æœ¬å†…å®¹
        text_content = ""
        if isinstance(slide_content, dict):
            # å¦‚æœcontentæ˜¯å­—å…¸ï¼Œå°è¯•æå–æ–‡æœ¬å­—æ®µ
            text_content = slide_content.get('text', '') or slide_content.get('content', '') or str(slide_content)
        else:
            text_content = str(slide_content)

        if not text_content or text_content.strip() == "":
            return [], []

        # æ„å»ºæç¤ºè¯
        prompt = f"""
è¯·ä»ä»¥ä¸‹Arduino/ç”µå­å·¥ç¨‹è¯¾ç¨‹å¹»ç¯ç‰‡å†…å®¹ä¸­æŠ½å–å®ä½“å’Œå…³ç³»ã€‚

å†…å®¹ï¼š{text_content}

è¯·è¯†åˆ«ä»¥ä¸‹ç±»å‹çš„å®ä½“ï¼š
1. ç¡¬ä»¶ç»„ä»¶ï¼šArduinoæ¿ã€ä¼ æ„Ÿå™¨ã€LEDã€ç”µé˜»ã€ç”µå®¹ç­‰
2. æŠ€æœ¯æ¦‚å¿µï¼šPWMã€ä¸²å£é€šä¿¡ã€æ•°å­—ä¿¡å·ã€æ¨¡æ‹Ÿä¿¡å·ç­‰
3. å‚æ•°æ•°å€¼ï¼šç”µå‹å€¼ã€ç”µé˜»å€¼ã€å¼•è„šå·ã€é¢‘ç‡ç­‰
4. æ“ä½œæ­¥éª¤ï¼šè¿æ¥ã€ç¼–ç¨‹ã€æµ‹è¯•ã€è°ƒè¯•ç­‰
5. ä»£ç æ¦‚å¿µï¼šå‡½æ•°ã€å˜é‡ã€åº“æ–‡ä»¶ç­‰

è¯·è¯†åˆ«å®ä½“é—´çš„å…³ç³»ï¼š
- ç»„æˆå…³ç³»ï¼šAåŒ…å«Bã€Aç”±Bç»„æˆ
- è¿æ¥å…³ç³»ï¼šAè¿æ¥åˆ°Bã€Aæ¥å…¥B
- æ§åˆ¶å…³ç³»ï¼šAæ§åˆ¶Bã€Aé©±åŠ¨B
- å‚æ•°å…³ç³»ï¼šAçš„å‚æ•°æ˜¯Bã€Aè®¾ç½®ä¸ºB
- åŠŸèƒ½å…³ç³»ï¼šAç”¨äºBã€Aå®ç°B

ä¸¥æ ¼æŒ‰ç…§ä»¥ä¸‹JSONæ ¼å¼è¿”å›ï¼Œä¸è¦æ·»åŠ ä»»ä½•å…¶ä»–å†…å®¹ï¼š
{{
    "entities": [
        {{"name": "å®ä½“åç§°", "type": "å®ä½“ç±»å‹", "description": "å®ä½“æè¿°"}}
    ],
    "relationships": [
        {{"source": "æºå®ä½“", "target": "ç›®æ ‡å®ä½“", "relation": "å…³ç³»ç±»å‹"}}
    ]
}}
"""

        try:
            response = self.deepseek.chat_completions_create([
                {"role": "user", "content": prompt}
            ])

            content = response['choices'][0]['message']['content']

            # æå–JSONéƒ¨åˆ†
            json_start = content.find('{')
            json_end = content.rfind('}') + 1

            if json_start != -1 and json_end != -1:
                json_str = content[json_start:json_end]
                result = json.loads(json_str)

                # æ·»åŠ slideä¿¡æ¯
                entities = result.get('entities', [])
                for entity in entities:
                    entity['slide'] = slide_num
                    entity['source'] = 'slide_text'

                relationships = result.get('relationships', [])
                for rel in relationships:
                    rel['slide'] = slide_num
                    rel['source'] = 'slide_text'

                return entities, relationships

        except Exception as e:
            print(f"   âš ï¸ å¹»ç¯ç‰‡ {slide_num} å®ä½“æŠ½å–å¤±è´¥: {e}")

        return [], []

    def _extract_from_image(self, image_data: Dict) -> List[Dict]:
        """ä»å›¾ç‰‡æ•°æ®ä¸­æŠ½å–å®ä½“"""
        entities = []
        slide_num = image_data.get('slide_number', 0)
        image_path = image_data.get('image_path', '')
        filename = image_data.get('filename', '')

        # 1. åŸºäºå›¾ç‰‡æè¿°æŠ½å–å®ä½“
        descriptions = image_data.get('descriptions', [])
        for desc_item in descriptions:
            desc_text = desc_item.get('description', '')
            confidence = desc_item.get('confidence', 0)

            if desc_text and confidence > 0.05:  # ç½®ä¿¡åº¦é˜ˆå€¼
                entities.append({
                    'name': desc_text,
                    'type': 'image_concept',
                    'description': f'ä»å›¾ç‰‡æè¿°ä¸­è¯†åˆ«: {desc_text}',
                    'confidence': confidence,
                    'source': 'image_description',
                    'slide': slide_num,
                    'image_path': image_path,
                    'filename': filename
                })

        # 2. åŸºäºOCRæ–‡æœ¬æŠ½å–å®ä½“ï¼ˆå¦‚æœæœ‰OCRæ–‡æœ¬ï¼‰
        ocr_text = image_data.get('ocr_text', '')
        if ocr_text:
            # Arduinoå…³é”®è¯åŒ¹é…
            for keyword in self.arduino_keywords:
                if keyword.lower() in ocr_text.lower():
                    entities.append({
                        'name': keyword,
                        'type': 'hardware_component',
                        'description': f'ä»å›¾ç‰‡OCRä¸­è¯†åˆ«çš„{keyword}',
                        'source': 'image_ocr',
                        'slide': slide_num,
                        'image_path': image_path,
                        'filename': filename
                    })

        # 3. åŸºäºæ–‡ä»¶åæŠ½å–å®ä½“ï¼ˆå¦‚æœæ–‡ä»¶ååŒ…å«æœ‰ç”¨ä¿¡æ¯ï¼‰
        if 'arduino' in filename.lower():
            entities.append({
                'name': 'Arduino',
                'type': 'hardware_platform',
                'description': 'ä»æ–‡ä»¶åè¯†åˆ«çš„Arduinoå¹³å°',
                'source': 'filename',
                'slide': slide_num,
                'image_path': image_path,
                'filename': filename
            })

        return entities

    def _deduplicate_entities(self, entities: List[Dict]) -> List[Dict]:
        """å®ä½“å»é‡"""
        seen = set()
        unique_entities = []

        for entity in entities:
            key = (entity['name'].lower(), entity['type'])
            if key not in seen:
                seen.add(key)
                unique_entities.append(entity)

        return unique_entities

    def _deduplicate_relationships(self, relationships: List[Dict]) -> List[Dict]:
        """å…³ç³»å»é‡"""
        seen = set()
        unique_relationships = []

        for rel in relationships:
            key = (rel['source'].lower(), rel['target'].lower(), rel['relation'])
            if key not in seen:
                seen.add(key)
                unique_relationships.append(rel)

        return unique_relationships


def extract_entities_from_output(output_dir: str, deepseek_api_key: str) -> ExtractedTriple:
    """ä»å¤šæ¨¡æ€è¾“å‡ºä¸­æŠ½å–å®ä½“å…³ç³»çš„ä¸»å‡½æ•°"""
    extractor = EntityExtractor(deepseek_api_key)

    # åŠ è½½æ•°æ®
    multimodal_data = extractor.load_multimodal_data(output_dir)

    # æŠ½å–å®ä½“å…³ç³»
    extracted_data = extractor.extract_entities_from_multimodal(multimodal_data)

    return extracted_data


import sys
import os

# è·å–å½“å‰æ–‡ä»¶çš„ç»å¯¹è·¯å¾„
current_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(current_dir)

# æ·»åŠ è·¯å¾„
sys.path.append(parent_dir)  # é¡¹ç›®æ ¹ç›®å½•
sys.path.append(current_dir)  # srcç›®å½•
import json
import fitz  # PyMuPDF
import torch
import numpy as np
from PIL import Image, ImageEnhance
from transformers import CLIPProcessor, CLIPModel
import datetime
import cv2
import easyocr
import re
import pdfplumber
import camelot
import csv

'''
try:
    # æ‡’åŠ è½½é«˜çº§PPTXå¤„ç†å™¨ï¼ˆè‹¥ä¸å¯ç”¨åˆ™å¿½ç•¥ï¼‰
    from advanced_pptx_processor import process_pptx_file_advanced
except Exception:
    process_pptx_file_advanced = None
'''


class MultimodalPreprocessor:
    def __init__(self):
        """åˆå§‹åŒ–å¤šæ¨¡æ€é¢„å¤„ç†å·¥å…·"""
        print("ğŸš€ å¼€å§‹åˆå§‹åŒ–å¤šæ¨¡æ€é¢„å¤„ç†å·¥å…·...")

        # æ£€æµ‹è®¾å¤‡
        print("ğŸ“± æ£€æµ‹è®¡ç®—è®¾å¤‡...")
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"âœ“ ä½¿ç”¨è®¾å¤‡: {self.device}")

        # åˆ›å»ºè¾“å‡ºç›®å½•
        print("ğŸ“ åˆ›å»ºè¾“å‡ºç›®å½•...")
        os.makedirs("output/text", exist_ok=True)
        os.makedirs("output/images", exist_ok=True)
        os.makedirs("output/formulas", exist_ok=True)
        os.makedirs("output/tables", exist_ok=True)
        os.makedirs("output/code", exist_ok=True)
        print("âœ“ è¾“å‡ºç›®å½•åˆ›å»ºå®Œæˆ")

        # åˆå§‹åŒ–CLIPæ¨¡å‹ï¼ˆå¯èƒ½è¾ƒæ…¢ï¼‰
        print("ğŸ¤– æ­£åœ¨åŠ è½½CLIPæ¨¡å‹...")
        print("   â³ æœ¬åœ°æ¨¡å‹åŠ è½½ä¸­ï¼Œè¯·ç¨å€™...")
        # ===== ä¿®æ”¹è¿™é‡Œï¼šä½¿ç”¨ä½ çš„æœ¬åœ°CLIPæ¨¡å‹è·¯å¾„ =====
        local_clip_path = r"F:\Models\clip-vit-base-patch32"
        try:
            if os.path.exists(local_clip_path):
                print(f"   ğŸ“ æ‰¾åˆ°æœ¬åœ°æ¨¡å‹: {local_clip_path}")
                self.clip_model = CLIPModel.from_pretrained(local_clip_path, local_files_only=True).to(self.device)
                print("   âœ“ æœ¬åœ°CLIPæ¨¡å‹åŠ è½½å®Œæˆ")
            else:
                print(f"   âŒ æœ¬åœ°æ¨¡å‹è·¯å¾„ä¸å­˜åœ¨: {local_clip_path}")
                raise FileNotFoundError("æœ¬åœ°æ¨¡å‹ä¸å­˜åœ¨")
        except Exception as e:
            print(f"   âŒ æœ¬åœ°CLIPæ¨¡å‹åŠ è½½å¤±è´¥: {e}")
            print("   â³ å°è¯•åœ¨çº¿ä¸‹è½½CLIPæ¨¡å‹ï¼Œè¿™å¯èƒ½éœ€è¦å‡ åˆ†é’Ÿ...")
            try:
                self.clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32").to(self.device)
                print("   âœ“ åœ¨çº¿CLIPæ¨¡å‹ä¸‹è½½å¹¶åŠ è½½å®Œæˆ")
            except Exception as e2:
                print(f"   âŒ CLIPæ¨¡å‹åŠ è½½å®Œå…¨å¤±è´¥: {e2}")
                raise e2

        print("ğŸ”§ æ­£åœ¨åŠ è½½CLIPå¤„ç†å™¨...")
        print("   â³ å¤„ç†å™¨åŠ è½½ä¸­ï¼ˆå¯èƒ½éœ€è¦ä¸‹è½½ï¼‰...")
        # ===== ä¿®æ”¹è¿™é‡Œï¼šä½¿ç”¨ä½ çš„æœ¬åœ°CLIPå¤„ç†å™¨è·¯å¾„ =====
        local_clip_path = r"F:\Models\clip-vit-base-patch32"
        try:
            if os.path.exists(local_clip_path):
                print(f"   ğŸ“ ä½¿ç”¨æœ¬åœ°å¤„ç†å™¨: {local_clip_path}")
                self.clip_processor = CLIPProcessor.from_pretrained(local_clip_path, local_files_only=True)
                print("   âœ“ æœ¬åœ°CLIPå¤„ç†å™¨åŠ è½½å®Œæˆ")
            else:
                print(f"   âŒ æœ¬åœ°å¤„ç†å™¨è·¯å¾„ä¸å­˜åœ¨ï¼Œä½¿ç”¨åœ¨çº¿ç‰ˆæœ¬")
                self.clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
                print("   âœ“ åœ¨çº¿CLIPå¤„ç†å™¨åŠ è½½å®Œæˆ")
        except Exception as e:
            print(f"   âŒ CLIPå¤„ç†å™¨åŠ è½½å¤±è´¥: {e}")
            raise e

        # åˆå§‹åŒ–OCRå¼•æ“ï¼ˆé¦–æ¬¡è¿è¡Œè¾ƒæ…¢ï¼‰
        print("ğŸ‘ æ­£åœ¨åˆå§‹åŒ–OCRå¼•æ“...")
        print("   â³ é¦–æ¬¡è¿è¡Œéœ€è¦ä¸‹è½½æ¨¡å‹æ–‡ä»¶ï¼Œè¿™å¯èƒ½éœ€è¦å‡ åˆ†é’Ÿï¼Œè¯·è€å¿ƒç­‰å¾…...")
        print("   ğŸ“¥ æ­£åœ¨ä¸‹è½½ä¸­æ–‡å’Œè‹±æ–‡OCRæ¨¡å‹...")
        try:
            self.ocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False)  # å¼ºåˆ¶ä½¿ç”¨CPUé¿å…GPUé—®é¢˜
            print("   âœ“ OCRå¼•æ“åˆå§‹åŒ–å®Œæˆ")
        except Exception as e:
            print(f"   âš  OCRåˆå§‹åŒ–å¤±è´¥ï¼Œå°†è·³è¿‡OCRåŠŸèƒ½: {e}")
            print("   å°†ç»§ç»­è¿è¡Œï¼Œä½†è·³è¿‡OCRå…¬å¼è¯†åˆ«åŠŸèƒ½")
            self.ocr_reader = None

        # å­˜å‚¨å¤„ç†ç»“æœ
        self.results = []

        print("ğŸ‰ å¤šæ¨¡æ€é¢„å¤„ç†å·¥å…·åˆå§‹åŒ–å®Œæˆï¼")
        print("=" * 50)

    def process_pdf(self, file_path):
        """å¤„ç†PDFæ–‡ä»¶ï¼Œæå–æ–‡æœ¬å’Œå›¾åƒ"""
        print(f"å¼€å§‹å¤„ç†PDFæ–‡ä»¶: {file_path}")
        doc = fitz.open(file_path)
        base_filename = os.path.splitext(os.path.basename(file_path))[0]

        # ä¸ºå½“å‰PDFæ–‡ä»¶å•ç‹¬è®°å½•ç»“æœ
        current_results = []
        original_results = self.results
        self.results = current_results

        try:
            for page_num in range(len(doc)):
                print(f"å¤„ç†ç¬¬ {page_num + 1}/{len(doc)} é¡µ...")
                page = doc.load_page(page_num)
                page_text = page.get_text()

                # å¤„ç†é¡µé¢å›¾åƒ
                image_list = page.get_images(full=True)
                page_images = []

                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]

                    # ä¿å­˜åŸå§‹å›¾åƒ
                    img_path = f"output/images/{base_filename}_p{page_num + 1}_img{img_index + 1}.{base_image['ext']}"
                    with open(img_path, "wb") as img_file:
                        img_file.write(image_bytes)

                    # å¤„ç†å›¾åƒå¹¶ä¿å­˜
                    image_data = self.process_image(img_path, page_text)
                    self.save_image_data(image_data, base_filename, page_num, img_index)
                    page_images.append(img_path)

                # å¤„ç†é¡µé¢æ–‡æœ¬
                text_data = self.process_text(page_text, page_num, page_images)
                text_data["source"] = f"{base_filename}_page{page_num + 1}"
                self.save_text_data(text_data, base_filename, page_num)

                # æå–é¡µé¢ä¸­çš„å…¬å¼ã€è¡¨æ ¼ã€ä»£ç 
                self.extract_formulas_from_page(page, page_text, base_filename, page_num)
                self.extract_tables_from_page(page, page_text, base_filename, page_num)
                self.extract_code_from_page(page_text, base_filename, page_num)

            # ä¿å­˜PDFä¸“ç”¨å…ƒæ•°æ®
            self.save_pdf_metadata(file_path, base_filename)
            print(f"PDFå¤„ç†å®Œæˆï¼ç»“æœä¿å­˜åœ¨output/{base_filename}_pdf_metadata.json")

        finally:
            # æ¢å¤åŸå§‹ç»“æœåˆ—è¡¨å¹¶åˆå¹¶å½“å‰ç»“æœ
            self.results = original_results
            self.results.extend(current_results)

    def process_text(self, text, page_num, page_images):
        """å¤„ç†æ–‡æœ¬å†…å®¹"""
        # æ¸…ç†æ–‡æœ¬
        cleaned_text = text.strip()
        if not cleaned_text:
            cleaned_text = "[é¡µé¢æ— æ–‡æœ¬å†…å®¹]"

        # ä½¿ç”¨CLIPç”Ÿæˆæ–‡æœ¬è¯­ä¹‰å‘é‡
        try:
            inputs = self.clip_processor(text=cleaned_text, return_tensors="pt", padding=True, truncation=True)
            inputs = {k: v.to(self.device) for k, v in inputs.items()}

            with torch.no_grad():
                text_features = self.clip_model.get_text_features(**inputs)
                text_vector = text_features.cpu().numpy()[0]
        except Exception as e:
            print(f"æ–‡æœ¬å‘é‡åŒ–å¤±è´¥: {e}")
            text_vector = np.zeros(512)  # CLIPé»˜è®¤å‘é‡ç»´åº¦

        return {
            "type": "text",
            "page": page_num + 1,
            "raw_text": cleaned_text,
            "word_count": len(cleaned_text),
            "associated_images": page_images,
            "text_vector": text_vector.tolist()
        }

    def process_image(self, image_path, page_text):
        """å¤„ç†å›¾åƒï¼ˆä½¿ç”¨CLIPï¼‰"""
        # å›¾åƒå¢å¼º
        enhanced_path = self.enhance_image(image_path)

        # è·å–å›¾åƒåŸºæœ¬ä¿¡æ¯
        try:
            with Image.open(image_path) as img:
                width, height = img.size
                format_type = img.format
                mode = img.mode
        except Exception as e:
            print(f"è¯»å–å›¾åƒä¿¡æ¯å¤±è´¥: {e}")
            width = height = 0
            format_type = mode = "unknown"

        # ä½¿ç”¨CLIPç”Ÿæˆå›¾åƒå‘é‡å’Œæè¿°
        try:
            image = Image.open(enhanced_path)
            inputs = self.clip_processor(images=image, return_tensors="pt").to(self.device)

            with torch.no_grad():
                image_features = self.clip_model.get_image_features(**inputs)
                image_vector = image_features.cpu().numpy()[0]

            # ç”Ÿæˆå›¾åƒæè¿°æ ‡ç­¾
            description_tags = self.generate_image_descriptions(enhanced_path)

        except Exception as e:
            print(f"å›¾åƒå¤„ç†å¤±è´¥: {e}")
            image_vector = np.zeros(512)
            description_tags = []

        return {
            "type": "image",
            "image_path": image_path,
            "enhanced_path": enhanced_path,
            "width": width,
            "height": height,
            "format": format_type,
            "mode": mode,
            "page_text_context": page_text[:200] + "..." if len(page_text) > 200 else page_text,
            "image_vector": image_vector.tolist(),
            "clip_descriptions": description_tags
        }

    def enhance_image(self, image_path):
        """å›¾åƒå¢å¼ºå¤„ç†"""
        img = Image.open(image_path)

        # å¯¹æ¯”åº¦å¢å¼º
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(1.5)

        # é”åº¦å¢å¼º
        enhancer = ImageEnhance.Sharpness(img)
        img = enhancer.enhance(2.0)

        # ä¿å­˜å¢å¼ºåçš„å›¾åƒ
        enhanced_path = image_path.replace(".", "_enhanced.")
        img.save(enhanced_path)

        return enhanced_path

    def clip_generate_description(self, image_path: str) -> str:
        """åŸºäºCLIPä¸ºå›¾ç‰‡ç”Ÿæˆæè¿°æ–‡æœ¬å¹¶ä¿å­˜ä¸ºJSONï¼Œè¿”å›æè¿°æ–‡ä»¶è·¯å¾„ã€‚"""
        try:
            descriptions = self.generate_image_descriptions(image_path)
        except Exception as e:
            print(f"ç”Ÿæˆå›¾ç‰‡æè¿°å¤±è´¥: {e}")
            descriptions = []

        base, _ = os.path.splitext(os.path.basename(image_path))
        desc_path = os.path.join("output", "images", f"{base}_desc.json")
        try:
            with open(desc_path, "w", encoding="utf-8") as f:
                json.dump({
                    "image_path": image_path,
                    "clip_descriptions": descriptions
                }, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"ä¿å­˜å›¾ç‰‡æè¿°å¤±è´¥: {e}")
        return desc_path

    def generate_image_descriptions(self, image_path):
        """ä½¿ç”¨CLIPç”Ÿæˆå›¾åƒæè¿°æ ‡ç­¾"""
        try:
            image = Image.open(image_path)
            image_input = self.clip_processor(images=image, return_tensors="pt").to(self.device)

            # é¢„å®šä¹‰çš„æè¿°å€™é€‰åˆ—è¡¨
            text_descriptions = [
                "ç§‘å­¦å›¾è¡¨", "æ•°å­¦å…¬å¼", "æ•°æ®å›¾è¡¨", "æµç¨‹å›¾",
                "å®éªŒè£…ç½®", "åˆ†å­ç»“æ„", "å‡ ä½•å›¾å½¢", "ç»Ÿè®¡å›¾è¡¨",
                "æŠ€æœ¯ç¤ºæ„å›¾", "æ¦‚å¿µå›¾", "ç½‘ç»œå›¾", "ç³»ç»Ÿæ¶æ„å›¾",
                "ç…§ç‰‡", "æ’å›¾", "ç¤ºä¾‹å›¾", "å¯¹æ¯”å›¾",
                "æ–‡æœ¬å›¾åƒ", "è¡¨æ ¼", "ä»£ç ", "æˆªå›¾"
            ]

            text_inputs = self.clip_processor(
                text=text_descriptions,
                return_tensors="pt",
                padding=True,
                truncation=True
            ).to(self.device)

            with torch.no_grad():
                image_features = self.clip_model.get_image_features(**image_input)
                text_features = self.clip_model.get_text_features(**text_inputs)

                # å½’ä¸€åŒ–ç‰¹å¾å‘é‡
                image_features = image_features / image_features.norm(dim=-1, keepdim=True)
                text_features = text_features / text_features.norm(dim=-1, keepdim=True)

                # è®¡ç®—ç›¸ä¼¼åº¦
                similarity = (image_features @ text_features.T).softmax(dim=-1)
                values, indices = similarity[0].topk(5)  # å–å‰5ä¸ªæœ€ç›¸ä¼¼çš„æè¿°

                descriptions = []
                for value, idx in zip(values.cpu().numpy(), indices.cpu().numpy()):
                    descriptions.append({
                        "description": text_descriptions[idx],
                        "confidence": float(value)
                    })

            return descriptions

        except Exception as e:
            print(f"å›¾åƒæè¿°ç”Ÿæˆé”™è¯¯: {image_path}, {str(e)}")
            return []

    def save_text_data(self, data, filename, page_num):
        """ä¿å­˜æ–‡æœ¬å¤„ç†ç»“æœ"""
        output_path = f"output/text/{filename}_p{page_num + 1}.json"
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        self.results.append({
            "type": "text",
            "page": page_num + 1,
            "file": output_path
        })

    def save_image_data(self, data, filename, page_num, img_index):
        """ä¿å­˜å›¾åƒå¤„ç†ç»“æœ"""
        output_path = f"output/images/{filename}_p{page_num + 1}_img{img_index + 1}.json"
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        self.results.append({
            "type": "image",
            "page": page_num + 1,
            "file": output_path
        })

    def save_pdf_metadata(self, file_path, filename):
        """ä¿å­˜PDFä¸“ç”¨å…ƒæ•°æ®æ–‡ä»¶"""
        # ç»Ÿè®¡ä¿¡æ¯
        text_files = [r for r in self.results if r["type"] == "text"]
        image_files = [r for r in self.results if r["type"] == "image"]
        formula_files = [r for r in self.results if r["type"] == "formula"]
        table_files = [r for r in self.results if r["type"] == "table"]
        code_files = [r for r in self.results if r["type"] == "code"]

        # è®¡ç®—è¡¨æ ¼ç»Ÿè®¡
        total_table_rows = sum(r.get("rows", 0) for r in table_files)
        total_table_columns = sum(r.get("columns", 0) for r in table_files)

        # è®¡ç®—ä»£ç ç»Ÿè®¡
        total_code_lines = sum(r.get("line_count", 0) for r in code_files)
        code_languages = list(set(r.get("language", "unknown") for r in code_files))

        metadata = {
            "processing_date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "source_file": filename,
            "source_path": file_path,
            "file_type": "PDF",
            "processing_method": "pymupdf_ocr_camelot",
            "statistics": {
                "total_pages": len(set(r["page"] for r in self.results)),
                "total_text_blocks": len(text_files),
                "total_images": len(image_files),
                "total_formulas": len(formula_files),
                "total_tables": len(table_files),
                "total_table_rows": total_table_rows,
                "total_table_columns": total_table_columns,
                "total_code_blocks": len(code_files),
                "total_code_lines": total_code_lines,
                "code_languages": code_languages
            },
            "files": {
                "text_files": text_files,
                "image_files": image_files,
                "formula_files": formula_files,
                "table_files": table_files,
                "code_files": code_files
            },
            "processing_info": {
                "clip_model": "openai/clip-vit-base-patch32",
                "device": self.device,
                "output_format": "JSON/CSV",
                "ocr_enabled": self.ocr_reader is not None,
                "extraction_features": ["text", "images", "formulas", "tables", "code"]
            }
        }

        with open(f"output/{filename}_pdf_metadata.json", "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)

    def extract_formulas_from_page(self, page, page_text, filename, page_num):
        """ä»é¡µé¢ä¸­æå–æ•°å­¦å…¬å¼"""
        formulas = []

        # 1. ä»æ–‡æœ¬ä¸­æå–LaTeXæ ¼å¼çš„å…¬å¼
        latex_patterns = [
            r'\$\$([^$]+)\$\$',  # å—çº§å…¬å¼ $$...$$
            r'\$([^$]+)\$',  # è¡Œå†…å…¬å¼ $...$
            r'\\begin\{equation\}(.*?)\\end\{equation\}',  # equationç¯å¢ƒ
            r'\\begin\{align\}(.*?)\\end\{align\}',  # alignç¯å¢ƒ
            r'\\begin\{math\}(.*?)\\end\{math\}',  # mathç¯å¢ƒ
        ]

        for i, pattern in enumerate(latex_patterns):
            matches = re.findall(pattern, page_text, re.DOTALL)
            for j, match in enumerate(matches):
                formula_data = {
                    "type": "formula",
                    "page": page_num + 1,
                    "formula_id": f"{filename}_p{page_num + 1}_formula{len(formulas) + 1}",
                    "content": match.strip(),
                    "format": "latex",
                    "extraction_method": f"pattern_{i + 1}",
                    "context": self.get_text_context(page_text, match, 100)
                }
                formulas.append(formula_data)

        # 2. ä»å›¾åƒä¸­è¯†åˆ«å…¬å¼ï¼ˆä½¿ç”¨OCRï¼‰- ç®€åŒ–ç‰ˆæœ¬é¿å…å¡ä½
        if self.ocr_reader and len(formulas) < 5:  # é™åˆ¶OCRå¤„ç†ï¼Œé¿å…å¡ä½
            try:
                # è·å–é¡µé¢å›¾åƒ
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                img_array = np.frombuffer(img_data, np.uint8)
                img = cv2.imdecode(img_array, cv2.IMREAD_COLOR)

                # OCRè¯†åˆ«ï¼Œè®¾ç½®è¶…æ—¶
                results = self.ocr_reader.readtext(img, width_ths=0.7, height_ths=0.7)

                for result in results[:3]:  # åªå¤„ç†å‰3ä¸ªç»“æœï¼Œé¿å…è¿‡å¤šå¤„ç†
                    text = result[1]
                    confidence = result[2]

                    # æ£€æŸ¥æ˜¯å¦åŒ…å«æ•°å­¦ç¬¦å·
                    math_symbols = ['âˆ‘', 'âˆ«', 'âˆ‚', 'âˆ†', 'âˆ‡', 'âˆ', 'Â±', 'â‰ ', 'â‰¤', 'â‰¥', 'Î±', 'Î²', 'Î³', 'Î´', 'Î¸', 'Î»', 'Î¼',
                                    'Ï€', 'Ïƒ', 'Ï†', 'Ïˆ', 'Ï‰']
                    if any(symbol in text for symbol in math_symbols) and confidence > 0.6:
                        if any(char.isdigit() or char in '+-*/=()[]{}^_' for char in text):
                            formula_data = {
                                "type": "formula",
                                "page": page_num + 1,
                                "formula_id": f"{filename}_p{page_num + 1}_formula{len(formulas) + 1}",
                                "content": text,
                                "format": "ocr_text",
                                "confidence": float(confidence),
                                "extraction_method": "ocr",
                                "bbox": [[float(pt[0]), float(pt[1])] for pt in result[0]]  # è½¬æ¢ä¸ºPythonåŸç”Ÿç±»å‹
                            }
                            formulas.append(formula_data)

            except Exception as e:
                print(f"OCRå…¬å¼è¯†åˆ«å¤±è´¥ (é¡µé¢ {page_num + 1}): {e}")

        # ä¿å­˜å…¬å¼æ•°æ®
        if formulas:
            self.save_formulas_data(formulas, filename, page_num)

        return formulas

    def extract_tables_from_page(self, page, page_text, filename, page_num):
        """ä»é¡µé¢ä¸­æå–è¡¨æ ¼"""
        tables = []

        try:
            # ä½¿ç”¨camelotæå–è¡¨æ ¼ - é™åˆ¶å¤„ç†æ—¶é—´ï¼Œé¿å…å¡ä½
            pdf_path = None
            for file in os.listdir("input"):
                if file.lower().endswith('.pdf') and filename in file:
                    pdf_path = os.path.join("input", file)
                    break

            if pdf_path and os.path.exists(pdf_path) and page_num < 10:  # åªå¤„ç†å‰10é¡µï¼Œé¿å…å¡ä½
                # æå–å½“å‰é¡µé¢çš„è¡¨æ ¼ï¼Œé™åˆ¶å¤„ç†
                camelot_tables = camelot.read_pdf(pdf_path, pages=str(page_num + 1), flavor='lattice')

                for i, table in enumerate(camelot_tables[:2]):  # åªå¤„ç†å‰2ä¸ªè¡¨æ ¼
                    if table.df is not None and not table.df.empty and len(table.df) > 0:
                        table_data = {
                            "type": "table",
                            "page": page_num + 1,
                            "table_id": f"{filename}_p{page_num + 1}_table{i + 1}",
                            "rows": len(table.df),
                            "columns": len(table.df.columns),
                            "data": table.df.to_dict('records'),
                            "extraction_method": "camelot",
                            "accuracy": getattr(table, 'accuracy', 0.0)
                        }
                        tables.append(table_data)

        except Exception as e:
            print(f"Camelotè¡¨æ ¼æå–è·³è¿‡ (é¡µé¢ {page_num + 1}): {e}")

        # å¤‡ç”¨æ–¹æ³•ï¼šä»æ–‡æœ¬ä¸­è¯†åˆ«è¡¨æ ¼æ¨¡å¼
        table_patterns = self.detect_text_tables(page_text)
        for i, pattern in enumerate(table_patterns):
            table_data = {
                "type": "table",
                "page": page_num + 1,
                "table_id": f"{filename}_p{page_num + 1}_texttable{i + 1}",
                "content": pattern,
                "extraction_method": "text_pattern",
                "context": self.get_text_context(page_text, pattern, 50)
            }
            tables.append(table_data)

        # ä¿å­˜è¡¨æ ¼æ•°æ®
        if tables:
            self.save_tables_data(tables, filename, page_num)

        return tables

    def extract_code_from_page(self, page_text, filename, page_num):
        """ä»é¡µé¢æ–‡æœ¬ä¸­æå–ä»£ç å—"""
        code_blocks = []

        # ä»£ç å—æ¨¡å¼
        code_patterns = [
            r'```(\w*)\n(.*?)```',  # Markdownä»£ç å—
            r'`([^`]+)`',  # è¡Œå†…ä»£ç 
            r'(?:^|\n)((?:    |\t)[^\n]+(?:\n(?:    |\t)[^\n]+)*)',  # ç¼©è¿›ä»£ç å—
        ]

        # ç¼–ç¨‹è¯­è¨€å…³é”®å­—
        programming_keywords = [
            'def ', 'class ', 'import ', 'from ', 'if ', 'else:', 'for ', 'while ', 'return',
            'function', 'var ', 'let ', 'const ', 'console.log', 'print(', 'println',
            'public ', 'private ', 'static ', 'void ', 'int ', 'String ', 'boolean',
            '#include', 'using namespace', 'int main(', 'printf(', 'cout <<'
        ]

        for i, pattern in enumerate(code_patterns):
            matches = re.findall(pattern, page_text, re.DOTALL | re.MULTILINE)

            for j, match in enumerate(matches):
                if isinstance(match, tuple):
                    language = match[0] if match[0] else "unknown"
                    content = match[1] if len(match) > 1 else match[0]
                else:
                    content = match
                    language = "unknown"

                # æ£€æŸ¥æ˜¯å¦åŒ…å«ç¼–ç¨‹å…³é”®å­—
                if any(keyword in content for keyword in programming_keywords) or len(content.strip()) > 20:
                    code_data = {
                        "type": "code",
                        "page": page_num + 1,
                        "code_id": f"{filename}_p{page_num + 1}_code{len(code_blocks) + 1}",
                        "content": content.strip(),
                        "language": language,
                        "extraction_method": f"pattern_{i + 1}",
                        "line_count": len(content.strip().split('\n')),
                        "context": self.get_text_context(page_text, content, 100)
                    }
                    code_blocks.append(code_data)

        # ä¿å­˜ä»£ç æ•°æ®
        if code_blocks:
            self.save_code_data(code_blocks, filename, page_num)

        return code_blocks

    def get_text_context(self, full_text, target_text, context_length=100):
        """è·å–ç›®æ ‡æ–‡æœ¬çš„ä¸Šä¸‹æ–‡"""
        try:
            index = full_text.find(target_text)
            if index == -1:
                return target_text

            start = max(0, index - context_length)
            end = min(len(full_text), index + len(target_text) + context_length)
            return full_text[start:end]
        except:
            return target_text

    def detect_text_tables(self, text):
        """ä»æ–‡æœ¬ä¸­æ£€æµ‹è¡¨æ ¼æ¨¡å¼"""
        tables = []
        lines = text.split('\n')

        # å¯»æ‰¾åŒ…å«å¤šä¸ªåˆ¶è¡¨ç¬¦æˆ–ç©ºæ ¼åˆ†éš”çš„è¡Œ
        table_lines = []
        for line in lines:
            # æ£€æŸ¥æ˜¯å¦åŒ…å«è¡¨æ ¼ç‰¹å¾ï¼šå¤šä¸ªåˆ¶è¡¨ç¬¦ã€ç«–çº¿åˆ†éš”ç¬¦ç­‰
            if '\t' in line and line.count('\t') >= 2:
                table_lines.append(line)
            elif '|' in line and line.count('|') >= 2:
                table_lines.append(line)
            elif re.search(r'\s{3,}', line) and len(line.split()) >= 3:
                table_lines.append(line)
            else:
                if table_lines and len(table_lines) >= 2:
                    tables.append('\n'.join(table_lines))
                table_lines = []

        # æ£€æŸ¥æœ€åä¸€ç»„
        if table_lines and len(table_lines) >= 2:
            tables.append('\n'.join(table_lines))

        return tables

    def save_formulas_data(self, formulas, filename, page_num):
        """ä¿å­˜å…¬å¼æ•°æ®"""
        # JSONæ ¼å¼ä¿å­˜
        json_path = f"output/formulas/{filename}_p{page_num + 1}_formulas.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(formulas, f, ensure_ascii=False, indent=2)

        # CSVæ ¼å¼ä¿å­˜
        csv_path = f"output/formulas/{filename}_p{page_num + 1}_formulas.csv"
        if formulas:
            df = pd.DataFrame(formulas)
            df.to_csv(csv_path, index=False, encoding="utf-8")

        # è®°å½•åˆ°ç»“æœ
        for formula in formulas:
            self.results.append({
                "type": "formula",
                "page": page_num + 1,
                "file": json_path,
                "formula_id": formula["formula_id"]
            })

    def save_tables_data(self, tables, filename, page_num):
        """ä¿å­˜è¡¨æ ¼æ•°æ®"""
        # JSONæ ¼å¼ä¿å­˜
        json_path = f"output/tables/{filename}_p{page_num + 1}_tables.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(tables, f, ensure_ascii=False, indent=2)

        # ä¸ºæ¯ä¸ªè¡¨æ ¼å•ç‹¬ä¿å­˜CSV
        for i, table in enumerate(tables):
            if table.get("data") and isinstance(table["data"], list):
                csv_path = f"output/tables/{table['table_id']}.csv"
                try:
                    df = pd.DataFrame(table["data"])
                    df.to_csv(csv_path, index=False, encoding="utf-8")
                except Exception as e:
                    print(f"ä¿å­˜è¡¨æ ¼CSVå¤±è´¥: {e}")

        # è®°å½•åˆ°ç»“æœ
        for table in tables:
            self.results.append({
                "type": "table",
                "page": page_num + 1,
                "file": json_path,
                "table_id": table["table_id"],
                "rows": table.get("rows", 0),
                "columns": table.get("columns", 0)
            })

    def save_code_data(self, code_blocks, filename, page_num):
        """ä¿å­˜ä»£ç æ•°æ®"""
        # JSONæ ¼å¼ä¿å­˜
        json_path = f"output/code/{filename}_p{page_num + 1}_code.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(code_blocks, f, ensure_ascii=False, indent=2)

        # CSVæ ¼å¼ä¿å­˜
        csv_path = f"output/code/{filename}_p{page_num + 1}_code.csv"
        if code_blocks:
            df = pd.DataFrame(code_blocks)
            df.to_csv(csv_path, index=False, encoding="utf-8")

        # ä¸ºæ¯ä¸ªä»£ç å—å•ç‹¬ä¿å­˜æ–‡ä»¶
        for code in code_blocks:
            if code.get("language") and code.get("language") != "unknown":
                ext = self.get_file_extension(code["language"])
                code_file_path = f"output/code/{code['code_id']}.{ext}"
                with open(code_file_path, "w", encoding="utf-8") as f:
                    f.write(code["content"])

        # è®°å½•åˆ°ç»“æœ
        for code in code_blocks:
            self.results.append({
                "type": "code",
                "page": page_num + 1,
                "file": json_path,
                "code_id": code["code_id"],
                "language": code.get("language", "unknown"),
                "line_count": code.get("line_count", 0)
            })

    def get_file_extension(self, language):
        """æ ¹æ®ç¼–ç¨‹è¯­è¨€è·å–æ–‡ä»¶æ‰©å±•å"""
        extensions = {
            "python": "py",
            "javascript": "js",
            "java": "java",
            "cpp": "cpp",
            "c": "c",
            "csharp": "cs",
            "php": "php",
            "ruby": "rb",
            "go": "go",
            "rust": "rs",
            "swift": "swift",
            "kotlin": "kt",
            "typescript": "ts",
            "html": "html",
            "css": "css",
            "sql": "sql",
            "shell": "sh",
            "bash": "sh",
            "powershell": "ps1"
        }
        return extensions.get(language.lower(), "txt")


import os
import json
import zipfile
import tempfile
import shutil
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import re
import subprocess


class AdvancedPPTProcessor:
    def __init__(self, preprocessor, fast_mode=False):
        """
        åˆå§‹åŒ–é«˜çº§PPTXå¤„ç†å™¨

        Args:
            preprocessor: MultimodalPreprocessorå®ä¾‹ï¼Œç”¨äºé‡ç”¨è¾“å‡ºç›®å½•ä¸ç»“æœè®°å½•
            fast_mode: å¿«é€Ÿæ¨¡å¼ï¼Œè·³è¿‡è€—æ—¶çš„CLIPæè¿°ç”Ÿæˆ
        """
        self.preprocessor = preprocessor
        self.fast_mode = fast_mode
        self.output_text_dir = "output/text"
        self.output_table_dir = "output/tables"
        self.output_img_dir = "output/images"

        # ç¡®ä¿è¾“å‡ºç›®å½•å­˜åœ¨
        os.makedirs(self.output_text_dir, exist_ok=True)
        os.makedirs(self.output_table_dir, exist_ok=True)
        os.makedirs(self.output_img_dir, exist_ok=True)

    def extract_all_images_via_zip(self, file_path):
        """
        é€šè¿‡ZIPè§£å‹å’ŒXMLè§£ææå–PPTXä¸­çš„æ‰€æœ‰å›¾ç‰‡

        Args:
            file_path: PPTXæ–‡ä»¶è·¯å¾„

        Returns:
            dict: åŒ…å«å¹»ç¯ç‰‡åˆ°å›¾ç‰‡æ˜ å°„å…³ç³»çš„å­—å…¸
        """
        print(f"å¼€å§‹é€šè¿‡ZIPæ–¹å¼æå–å›¾ç‰‡: {file_path}")

        base_filename = os.path.splitext(os.path.basename(file_path))[0]
        slide_image_mapping = {}

        # åˆ›å»ºä¸´æ—¶ç›®å½•
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # 1. è§£å‹PPTXæ–‡ä»¶
                print("æ­£åœ¨è§£å‹PPTXæ–‡ä»¶...")
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)

                # 2. æ‰¾åˆ°åª’ä½“ç›®å½•
                media_dir = os.path.join(temp_dir, "ppt", "media")
                slides_dir = os.path.join(temp_dir, "ppt", "slides")
                rels_dir = os.path.join(temp_dir, "ppt", "slides", "_rels")

                if not os.path.exists(media_dir):
                    print("æœªæ‰¾åˆ°mediaç›®å½•ï¼Œå¯èƒ½æ²¡æœ‰å›¾ç‰‡")
                    return slide_image_mapping

                print(f"æ‰¾åˆ°mediaç›®å½•: {media_dir}")
                print(f"åª’ä½“æ–‡ä»¶: {os.listdir(media_dir)}")

                # 3. éå†æ‰€æœ‰å¹»ç¯ç‰‡XMLæ–‡ä»¶
                if os.path.exists(slides_dir):
                    for slide_file in os.listdir(slides_dir):
                        if slide_file.startswith("slide") and slide_file.endswith(".xml"):
                            slide_num = self._extract_slide_number(slide_file)
                            if slide_num is None:
                                continue

                            print(f"å¤„ç†å¹»ç¯ç‰‡ {slide_num}: {slide_file}")

                            # è§£æå¹»ç¯ç‰‡XMLè·å–å›¾ç‰‡å…³ç³»ID
                            slide_xml_path = os.path.join(slides_dir, slide_file)
                            image_rids = self._parse_slide_xml_for_images(slide_xml_path)

                            if image_rids:
                                print(f"å¹»ç¯ç‰‡ {slide_num} ä¸­æ‰¾åˆ°å›¾ç‰‡å…³ç³»ID: {image_rids}")

                                # è§£æå…³ç³»æ–‡ä»¶è·å–å®é™…æ–‡ä»¶å
                                rels_file = slide_file + ".rels"
                                rels_path = os.path.join(rels_dir, rels_file)

                                if os.path.exists(rels_path):
                                    image_files = self._parse_rels_file(rels_path, image_rids)

                                    if image_files:
                                        slide_image_mapping[slide_num] = image_files
                                        print(f"å¹»ç¯ç‰‡ {slide_num} æ˜ å°„åˆ°å›¾ç‰‡: {image_files}")

                                        # å¤åˆ¶å›¾ç‰‡åˆ°è¾“å‡ºç›®å½•
                                        self._copy_images_to_output(media_dir, image_files,
                                                                    base_filename, slide_num)

                print(f"å›¾ç‰‡æå–å®Œæˆï¼Œæ˜ å°„å…³ç³»: {slide_image_mapping}")

            except Exception as e:
                print(f"ZIPæ–¹å¼å›¾ç‰‡æå–å¤±è´¥: {e}")
                import traceback
                traceback.print_exc()

        return slide_image_mapping

    def _extract_slide_number(self, slide_filename):
        """ä»å¹»ç¯ç‰‡æ–‡ä»¶åä¸­æå–ç¼–å·"""
        try:
            # slide1.xml -> 1
            import re
            match = re.search(r'slide(\d+)\.xml', slide_filename)
            if match:
                return int(match.group(1))
        except Exception:
            pass
        return None

    def _parse_slide_xml_for_images(self, slide_xml_path):
        """
        è§£æå¹»ç¯ç‰‡XMLæ–‡ä»¶ï¼ŒæŸ¥æ‰¾å›¾ç‰‡å¼•ç”¨

        Args:
            slide_xml_path: å¹»ç¯ç‰‡XMLæ–‡ä»¶è·¯å¾„

        Returns:
            list: å›¾ç‰‡å…³ç³»IDåˆ—è¡¨
        """
        image_rids = []

        try:
            tree = ET.parse(slide_xml_path)
            root = tree.getroot()

            # å®šä¹‰å‘½åç©ºé—´
            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
            }

            # æŸ¥æ‰¾æ‰€æœ‰a:blipå…ƒç´ ï¼ˆå›¾ç‰‡å¼•ç”¨ï¼‰
            blip_elements = root.findall('.//a:blip', namespaces)

            for blip in blip_elements:
                embed_attr = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if embed_attr:
                    image_rids.append(embed_attr)
                    print(f"æ‰¾åˆ°å›¾ç‰‡å¼•ç”¨ID: {embed_attr}")

        except Exception as e:
            print(f"è§£æå¹»ç¯ç‰‡XMLå¤±è´¥ {slide_xml_path}: {e}")

        return image_rids

    def _parse_rels_file(self, rels_path, image_rids):
        """
        è§£æå…³ç³»æ–‡ä»¶ï¼Œè·å–å…³ç³»IDåˆ°æ–‡ä»¶åçš„æ˜ å°„

        Args:
            rels_path: å…³ç³»æ–‡ä»¶è·¯å¾„
            image_rids: å›¾ç‰‡å…³ç³»IDåˆ—è¡¨

        Returns:
            list: å¯¹åº”çš„å›¾ç‰‡æ–‡ä»¶ååˆ—è¡¨
        """
        image_files = []

        try:
            tree = ET.parse(rels_path)
            root = tree.getroot()

            # å®šä¹‰å‘½åç©ºé—´
            namespaces = {
                'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'
            }

            # æŸ¥æ‰¾æ‰€æœ‰å…³ç³»
            for relationship in root.findall('.//rel:Relationship', namespaces):
                rel_id = relationship.get('Id')
                target = relationship.get('Target')
                rel_type = relationship.get('Type')

                # æ£€æŸ¥æ˜¯å¦æ˜¯å›¾ç‰‡å…³ç³»
                if (rel_id in image_rids and
                        target and
                        rel_type and
                        'image' in rel_type.lower()):
                    # æå–æ–‡ä»¶å (../media/image1.png -> image1.png)
                    filename = os.path.basename(target)
                    image_files.append(filename)
                    print(f"å…³ç³»æ˜ å°„: {rel_id} -> {filename}")

        except Exception as e:
            print(f"è§£æå…³ç³»æ–‡ä»¶å¤±è´¥ {rels_path}: {e}")

        return image_files

    def _copy_images_to_output(self, media_dir, image_files, base_filename, slide_num):
        """
        å°†å›¾ç‰‡å¤åˆ¶åˆ°è¾“å‡ºç›®å½•å¹¶ç”Ÿæˆæè¿°

        Args:
            media_dir: åª’ä½“æ–‡ä»¶æºç›®å½•
            image_files: å›¾ç‰‡æ–‡ä»¶ååˆ—è¡¨
            base_filename: åŸºç¡€æ–‡ä»¶å
            slide_num: å¹»ç¯ç‰‡ç¼–å·
        """
        for idx, image_file in enumerate(image_files, 1):
            try:
                source_path = os.path.join(media_dir, image_file)

                if os.path.exists(source_path):
                    # ç”Ÿæˆè¾“å‡ºæ–‡ä»¶å
                    file_ext = os.path.splitext(image_file)[1]
                    output_filename = f"{base_filename}_slide_{slide_num}_img_{idx}_zip{file_ext}"
                    output_path = os.path.join(self.output_img_dir, output_filename)

                    # å¤åˆ¶å›¾ç‰‡
                    shutil.copy2(source_path, output_path)
                    print(f"å¤åˆ¶å›¾ç‰‡: {source_path} -> {output_path}")

                    # ç”ŸæˆCLIPæè¿°ï¼ˆæ ¹æ®æ¨¡å¼å†³å®šæ˜¯å¦ç”Ÿæˆï¼‰
                    desc_path = None
                    if not self.fast_mode:
                        try:
                            desc_path = self.preprocessor.clip_generate_description(output_path)
                        except Exception as e:
                            print(f"ç”ŸæˆCLIPæè¿°å¤±è´¥ï¼Œè·³è¿‡: {e}")
                    else:
                        print("å¿«é€Ÿæ¨¡å¼ï¼šè·³è¿‡CLIPæè¿°ç”Ÿæˆ")

                    # è®°å½•åˆ°ç»“æœä¸­
                    self.preprocessor.results.append({
                        "type": "ppt_image_zip",
                        "page": slide_num,
                        "file": output_path,
                        "description_file": desc_path,
                        "extraction_method": "zip_xml_parsing",
                        "original_filename": image_file
                    })

                else:
                    print(f"æºå›¾ç‰‡æ–‡ä»¶ä¸å­˜åœ¨: {source_path}")

            except Exception as e:
                print(f"å¤åˆ¶å›¾ç‰‡å¤±è´¥ {image_file}: {e}")

    def extract_and_convert_equations(self, slide, slide_number):
        """
        å¤„ç†å¹»ç¯ç‰‡ä¸­çš„å…¬å¼

        Args:
            slide: python-pptxçš„Slideå¯¹è±¡
            slide_number: å¹»ç¯ç‰‡ç¼–å·

        Returns:
            list: åŒ…å«å…¬å¼ä¿¡æ¯çš„åˆ—è¡¨
        """
        equations = []

        for shape_index, shape in enumerate(slide.shapes):
            try:
                # æ£€æŸ¥å½¢çŠ¶æ˜¯å¦åŒ…å«æ–‡æœ¬æ¡†
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # è·å–å½¢çŠ¶çš„XMLå†…å®¹
                    shape_xml = self._get_shape_xml(shape)
                    if shape_xml:
                        # æ£€æŸ¥æ˜¯å¦åŒ…å«OMMLå…¬å¼æ ‡ç­¾
                        omml_content = self._extract_omml_from_xml(shape_xml)
                        if omml_content:
                            print(f"åœ¨å¹»ç¯ç‰‡ {slide_number} å½¢çŠ¶ {shape_index} ä¸­å‘ç°OMMLå…¬å¼")

                            # å°è¯•è½¬æ¢OMMLåˆ°LaTeX
                            latex_content = self._convert_omml_to_latex(omml_content)

                            equation_info = {
                                "slide_number": slide_number,
                                "shape_index": shape_index,
                                "type": "omml_formula",
                                "original_omml": omml_content[:500] + "..." if len(
                                    omml_content) > 500 else omml_content,
                                "latex": latex_content,
                                "conversion_success": latex_content is not None
                            }

                            equations.append(equation_info)

                            # æ·»åŠ åˆ°ç»“æœä¸­
                            self.preprocessor.results.append({
                                "type": "formula",
                                "page": slide_number,
                                "formula_type": "omml",
                                "latex": latex_content,
                                "source": f"slide_{slide_number}_shape_{shape_index}",
                                "conversion_method": "omml_to_latex"
                            })

                            continue

                # å¦‚æœæ²¡æœ‰æ‰¾åˆ°OMMLï¼Œæ£€æŸ¥æ˜¯å¦ä¸ºå¯èƒ½çš„å…¬å¼å›¾ç‰‡
                if self._is_potential_formula_image(shape):
                    print(f"åœ¨å¹»ç¯ç‰‡ {slide_number} å½¢çŠ¶ {shape_index} ä¸­å‘ç°æ½œåœ¨å…¬å¼å›¾ç‰‡")

                    # ä½¿ç”¨å›¾ç‰‡å¤„ç†æµç¨‹å¤„ç†å…¬å¼å›¾ç‰‡
                    formula_image_path = self._process_formula_image(shape, slide_number, shape_index)

                    if formula_image_path:
                        equation_info = {
                            "slide_number": slide_number,
                            "shape_index": shape_index,
                            "type": "formula_image",
                            "image_path": formula_image_path,
                            "latex": None,
                            "conversion_success": False
                        }

                        equations.append(equation_info)

                        # æ·»åŠ åˆ°ç»“æœä¸­
                        self.preprocessor.results.append({
                            "type": "formula",
                            "page": slide_number,
                            "formula_type": "image",
                            "image_path": formula_image_path,
                            "source": f"slide_{slide_number}_shape_{shape_index}",
                            "conversion_method": "image_fallback"
                        })

            except Exception as e:
                print(f"å¤„ç†å¹»ç¯ç‰‡ {slide_number} å½¢çŠ¶ {shape_index} æ—¶å‡ºé”™: {e}")
                continue

        return equations

    def _get_shape_xml(self, shape):
        """è·å–å½¢çŠ¶çš„XMLå†…å®¹"""
        try:
            # å°è¯•è·å–å½¢çŠ¶çš„å†…éƒ¨XML
            if hasattr(shape, '_element'):
                return ET.tostring(shape._element, encoding='unicode')
        except Exception as e:
            print(f"è·å–å½¢çŠ¶XMLå¤±è´¥: {e}")
        return None

    def _extract_omml_from_xml(self, xml_string):
        """ä»XMLä¸­æå–OMMLå†…å®¹"""
        try:
            # æŸ¥æ‰¾OMMLæ•°å­¦æ ‡ç­¾
            omml_patterns = [
                r'<m:oMath[^>]*>.*?</m:oMath>',
                r'<m:oMathPara[^>]*>.*?</m:oMathPara>',
                r'<math[^>]*>.*?</math>'  # ä¹Ÿæ£€æŸ¥æ ‡å‡†MathML
            ]

            for pattern in omml_patterns:
                matches = re.findall(pattern, xml_string, re.DOTALL | re.IGNORECASE)
                if matches:
                    return matches[0]

        except Exception as e:
            print(f"æå–OMMLå¤±è´¥: {e}")
        return None

    def _convert_omml_to_latex(self, omml_content):
        """å°†OMMLè½¬æ¢ä¸ºLaTeX"""
        try:
            # æ–¹æ³•1: å°è¯•ä½¿ç”¨pandoc
            latex_result = self._convert_via_pandoc(omml_content)
            if latex_result:
                return latex_result

            # æ–¹æ³•2: ç®€å•çš„æ–‡æœ¬æ›¿æ¢ä½œä¸ºå¤‡é€‰æ–¹æ¡ˆ
            latex_result = self._simple_omml_to_latex(omml_content)
            if latex_result:
                return latex_result

        except Exception as e:
            print(f"OMMLè½¬LaTeXå¤±è´¥: {e}")

        return None

    def _convert_via_pandoc(self, omml_content):
        """ä½¿ç”¨pandocè½¬æ¢OMMLåˆ°LaTeX"""
        try:
            # æ£€æŸ¥pandocæ˜¯å¦å¯ç”¨
            subprocess.run(['pandoc', '--version'],
                           capture_output=True, check=True)

            # åˆ›å»ºä¸´æ—¶æ–‡ä»¶
            with tempfile.NamedTemporaryFile(mode='w', suffix='.xml', delete=False) as temp_file:
                temp_file.write(f'<root>{omml_content}</root>')
                temp_file_path = temp_file.name

            try:
                # ä½¿ç”¨pandocè½¬æ¢
                result = subprocess.run([
                    'pandoc',
                    '-f', 'docx',
                    '-t', 'latex',
                    temp_file_path
                ], capture_output=True, text=True, check=True)

                return result.stdout.strip()

            finally:
                os.unlink(temp_file_path)

        except (subprocess.CalledProcessError, FileNotFoundError):
            print("Pandocä¸å¯ç”¨ï¼Œè·³è¿‡pandocè½¬æ¢")
        except Exception as e:
            print(f"Pandocè½¬æ¢å¤±è´¥: {e}")

        return None

    def _simple_omml_to_latex(self, omml_content):
        """ç®€å•çš„OMMLåˆ°LaTeXè½¬æ¢ï¼ˆåŸºæœ¬æ–‡æœ¬æ›¿æ¢ï¼‰"""
        try:
            # ç§»é™¤XMLæ ‡ç­¾ï¼Œæå–çº¯æ–‡æœ¬
            text_content = re.sub(r'<[^>]+>', '', omml_content)
            text_content = text_content.strip()

            if not text_content:
                return None

            # åŸºæœ¬çš„æ•°å­¦ç¬¦å·æ›¿æ¢
            replacements = {
                'â‰ˆ': r'\approx',
                'â‰ ': r'\neq',
                'â‰¤': r'\leq',
                'â‰¥': r'\geq',
                'âˆ': r'\infty',
                'Î±': r'\alpha',
                'Î²': r'\beta',
                'Î³': r'\gamma',
                'Î´': r'\delta',
                'Î¸': r'\theta',
                'Î»': r'\lambda',
                'Î¼': r'\mu',
                'Ï€': r'\pi',
                'Ïƒ': r'\sigma',
                'Ï†': r'\phi',
                'Ï‰': r'\omega',
                'âˆ‘': r'\sum',
                'âˆ«': r'\int',
                'âˆš': r'\sqrt',
                'Â±': r'\pm',
                'Ã—': r'\times',
                'Ã·': r'\div'
            }

            for symbol, latex in replacements.items():
                text_content = text_content.replace(symbol, latex)

            # åŒ…è£…åœ¨æ•°å­¦ç¯å¢ƒä¸­
            return f"${text_content}$"

        except Exception as e:
            print(f"ç®€å•è½¬æ¢å¤±è´¥: {e}")

        return None

    def _is_potential_formula_image(self, shape):
        """åˆ¤æ–­å½¢çŠ¶æ˜¯å¦å¯èƒ½æ˜¯å…¬å¼å›¾ç‰‡"""
        try:
            # æ£€æŸ¥æ˜¯å¦ä¸ºå›¾ç‰‡ç±»å‹
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True

            # æ£€æŸ¥æ˜¯å¦ä¸ºåŒ…å«å¤æ‚è·¯å¾„çš„å½¢çŠ¶ï¼ˆå¯èƒ½æ˜¯çŸ¢é‡å…¬å¼ï¼‰
            if shape.shape_type in [MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.AUTO_SHAPE]:
                return True

            # æ£€æŸ¥å½¢çŠ¶å¤§å°ï¼ˆå°çš„å½¢çŠ¶å¯èƒ½æ˜¯å…¬å¼ï¼‰
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                # å‡è®¾å…¬å¼é€šå¸¸æ¯”è¾ƒå°ï¼ˆå®½åº¦å’Œé«˜åº¦éƒ½å°äºæŸä¸ªé˜ˆå€¼ï¼‰
                max_formula_size = 200000  # EMUå•ä½
                if shape.width < max_formula_size and shape.height < max_formula_size:
                    return True

        except Exception as e:
            print(f"æ£€æŸ¥æ½œåœ¨å…¬å¼å›¾ç‰‡å¤±è´¥: {e}")

        return False

    def _process_formula_image(self, shape, slide_number, shape_index):
        """å¤„ç†å…¬å¼å›¾ç‰‡"""
        try:
            # å¦‚æœæ˜¯å›¾ç‰‡ç±»å‹ï¼Œå°è¯•å¯¼å‡ºå›¾ç‰‡
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_filename = f"formula_slide_{slide_number}_shape_{shape_index}.png"
                image_path = os.path.join(self.output_img_dir, image_filename)

                # è¿™é‡Œéœ€è¦å®ç°å›¾ç‰‡å¯¼å‡ºé€»è¾‘
                # ç”±äºpython-pptxçš„é™åˆ¶ï¼Œå¯èƒ½éœ€è¦ä½¿ç”¨å…¶ä»–æ–¹æ³•
                print(f"è¯†åˆ«åˆ°å…¬å¼å›¾ç‰‡ï¼Œä½†éœ€è¦é¢å¤–çš„å¯¼å‡ºé€»è¾‘: {image_filename}")

                return image_path

        except Exception as e:
            print(f"å¤„ç†å…¬å¼å›¾ç‰‡å¤±è´¥: {e}")

        return None

    def process_pptx_file_advanced(self, file_path):
        """
        é«˜çº§PPTXå¤„ç†ï¼šç»“åˆä¼ ç»Ÿæ–¹æ³•å’ŒZIPè§£æ

        Args:
            file_path: PPTXæ–‡ä»¶è·¯å¾„
        """
        print(f"å¼€å§‹é«˜çº§PPTXå¤„ç†: {file_path}")
        base_filename = os.path.splitext(os.path.basename(file_path))[0]

        # é‡ç½®ç»“æœè®°å½•ï¼Œä¸ºå½“å‰PPTXæ–‡ä»¶å•ç‹¬è®°å½•
        current_results = []
        original_results = self.preprocessor.results
        self.preprocessor.results = current_results

        try:
            # 1. ä½¿ç”¨ä¼ ç»Ÿpython-pptxæ–¹æ³•å¤„ç†æ–‡æœ¬å’Œè¡¨æ ¼
            self._process_text_and_tables_traditional(file_path, base_filename)

            # 2. ä½¿ç”¨ZIPæ–¹æ³•æå–æ‰€æœ‰å›¾ç‰‡
            slide_image_mapping = self.extract_all_images_via_zip(file_path)

            # 3. ç”ŸæˆPPTXä¸“ç”¨å…ƒæ•°æ®
            self._save_pptx_metadata(file_path, base_filename, slide_image_mapping)

            print(f"é«˜çº§PPTXå¤„ç†å®Œæˆ: {file_path}")
            print(f"PPTXå…ƒæ•°æ®å·²ä¿å­˜: output/{base_filename}_pptx_metadata.json")

        except Exception as e:
            print(f"é«˜çº§PPTXå¤„ç†å¤±è´¥: {e}")
            import traceback
            traceback.print_exc()
        finally:
            # æ¢å¤åŸå§‹ç»“æœåˆ—è¡¨å¹¶åˆå¹¶å½“å‰ç»“æœ
            self.preprocessor.results = original_results
            self.preprocessor.results.extend(current_results)

    def _process_text_and_tables_traditional(self, file_path, base_filename):
        """ä½¿ç”¨ä¼ ç»Ÿpython-pptxæ–¹æ³•å¤„ç†æ–‡æœ¬å’Œè¡¨æ ¼"""
        prs = Presentation(file_path)

        for slide_index, slide in enumerate(prs.slides, start=1):
            # å¤„ç†å…¬å¼
            equations = self.extract_and_convert_equations(slide, slide_index)
            if equations:
                print(f"åœ¨å¹»ç¯ç‰‡ {slide_index} ä¸­æ‰¾åˆ° {len(equations)} ä¸ªå…¬å¼")

                # ä¿å­˜å…¬å¼ä¿¡æ¯åˆ°JSONæ–‡ä»¶
                formulas_json_path = f"output/formulas/{base_filename}_slide_{slide_index}_formulas.json"
                os.makedirs("output/formulas", exist_ok=True)

                formulas_output = {
                    "slide_number": slide_index,
                    "source_file": base_filename,
                    "equations_count": len(equations),
                    "equations": equations,
                    "processing_date": str(pd.Timestamp.now())
                }

                with open(formulas_json_path, "w", encoding="utf-8") as f:
                    json.dump(formulas_output, f, ensure_ascii=False, indent=2)

                print(f"å…¬å¼ä¿¡æ¯å·²ä¿å­˜åˆ°: {formulas_json_path}")

            # æå–æ–‡æœ¬
            slide_text_items = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text_content = []
                    for paragraph in shape.text_frame.paragraphs:
                        runs_text = ''.join(run.text for run in paragraph.runs)
                        text_content.append(runs_text if runs_text else paragraph.text)
                    final_text = "\n".join([t for t in text_content if t is not None])
                    if final_text.strip():
                        slide_text_items.append(final_text.strip())

            if slide_text_items:
                text_output = {
                    "type": "ppt_text",
                    "page": slide_index,
                    "source": f"{base_filename}_slide_{slide_index}",
                    "raw_text": "\n\n".join(slide_text_items)
                }
                text_json_path = f"{self.output_text_dir}/{base_filename}_slide_{slide_index}.json"
                with open(text_json_path, "w", encoding="utf-8") as f:
                    json.dump(text_output, f, ensure_ascii=False, indent=2)
                self.preprocessor.results.append({
                    "type": "ppt_text",
                    "page": slide_index,
                    "file": text_json_path
                })

            # æå–è¡¨æ ¼ï¼ˆä¼˜åŒ–ç‰ˆæœ¬ï¼‰
            table_counter = 0
            for shape in slide.shapes:
                if hasattr(shape, "has_table") and shape.has_table:
                    table_counter += 1
                    table = shape.table

                    # è·å–è¡¨æ ¼ä½ç½®ä¿¡æ¯ï¼ˆå¢å¼ºåŠŸèƒ½ï¼‰
                    table_position = {
                        "left": float(shape.left.inches) if shape.left else 0,
                        "top": float(shape.top.inches) if shape.top else 0,
                        "width": float(shape.width.inches) if shape.width else 0,
                        "height": float(shape.height.inches) if shape.height else 0
                    }

                    # æå–è¡¨æ ¼æ•°æ®
                    data_matrix = []
                    for row in table.rows:
                        row_values = []
                        for cell in row.cells:
                            # ä¼˜åŒ–ï¼šä½¿ç”¨ text_frame.text è·å–çº¯æ–‡æœ¬
                            try:
                                if cell.text_frame and cell.text_frame.text:
                                    cell_text = cell.text_frame.text.strip()
                                else:
                                    cell_text = cell.text.strip() if cell.text else ""
                            except Exception as e:
                                print(f"æå–å•å…ƒæ ¼æ–‡æœ¬å¤±è´¥: {e}")
                                cell_text = ""
                            row_values.append(cell_text)
                        data_matrix.append(row_values)

                    # æ£€æŸ¥æ˜¯å¦æœ‰æœ‰æ•ˆæ•°æ®
                    if data_matrix and any(any(cell for cell in row) for row in data_matrix):
                        # ä½¿ç”¨pandas DataFrameä¿å­˜ä¸ºCSV
                        df = pd.DataFrame(data_matrix)
                        csv_path = f"{self.output_table_dir}/{base_filename}_slide_{slide_index}_table_{table_counter}.csv"
                        df.to_csv(csv_path, index=False, header=False, encoding="utf-8")

                        # åˆ›å»ºè¡¨æ ¼JSONå…ƒæ•°æ®æ–‡ä»¶
                        table_metadata = {
                            "type": "ppt_table",
                            "source": f"{base_filename}_slide_{slide_index}",
                            "slide_number": slide_index,
                            "table_index": table_counter,
                            "dimensions": {
                                "rows": len(data_matrix),
                                "columns": len(data_matrix[0]) if data_matrix else 0
                            },
                            "position": table_position,
                            "data_preview": data_matrix[:3] if len(data_matrix) > 0 else [],  # å‰3è¡Œé¢„è§ˆ
                            "csv_file": csv_path
                        }

                        # ä¿å­˜è¡¨æ ¼å…ƒæ•°æ®JSON
                        json_path = f"{self.output_table_dir}/{base_filename}_slide_{slide_index}_table_{table_counter}.json"
                        with open(json_path, "w", encoding="utf-8") as f:
                            json.dump(table_metadata, f, ensure_ascii=False, indent=2)

                        # è®°å½•åˆ°ä¸»ç»“æœä¸­ï¼ˆå¢å¼ºå…ƒæ•°æ®ï¼‰
                        self.preprocessor.results.append({
                            "type": "ppt_table",
                            "page": slide_index,
                            "table_index": table_counter,
                            "file": csv_path,
                            "metadata_file": json_path,
                            "dimensions": {
                                "rows": len(data_matrix),
                                "columns": len(data_matrix[0]) if data_matrix else 0
                            },
                            "position": table_position,
                            "extraction_method": "python_pptx_optimized"
                        })

                        print(
                            f"âœ“ æå–è¡¨æ ¼ {table_counter}: {len(data_matrix)}è¡Œ x {len(data_matrix[0]) if data_matrix else 0}åˆ—")
                        print(f"  ä½ç½®: left={table_position['left']:.2f}in, top={table_position['top']:.2f}in")
                    else:
                        print(f"âš  è·³è¿‡ç©ºè¡¨æ ¼ {table_counter}")

    def _save_pptx_metadata(self, file_path, base_filename, slide_image_mapping):
        """
        ä¿å­˜PPTXä¸“ç”¨å…ƒæ•°æ®æ–‡ä»¶

        Args:
            file_path: åŸå§‹PPTXæ–‡ä»¶è·¯å¾„
            base_filename: åŸºç¡€æ–‡ä»¶å
            slide_image_mapping: å¹»ç¯ç‰‡åˆ°å›¾ç‰‡çš„æ˜ å°„å…³ç³»
        """
        import datetime

        # ç»Ÿè®¡å„ç±»å‹æ–‡ä»¶
        text_files = [r for r in self.preprocessor.results if r["type"] == "ppt_text"]
        table_files = [r for r in self.preprocessor.results if r["type"] == "ppt_table"]
        image_files_traditional = [r for r in self.preprocessor.results if r["type"] == "ppt_image"]
        image_files_zip = [r for r in self.preprocessor.results if r["type"] == "ppt_image_zip"]

        # è®¡ç®—å¹»ç¯ç‰‡ç»Ÿè®¡
        total_slides = len(set(r["page"] for r in self.preprocessor.results if "page" in r))
        slides_with_images = len(slide_image_mapping)
        total_images_zip = sum(len(images) for images in slide_image_mapping.values())

        # è®¡ç®—è¡¨æ ¼ç»Ÿè®¡
        total_tables = len(table_files)
        total_table_rows = sum(r.get("dimensions", {}).get("rows", 0) for r in table_files)
        total_table_columns = sum(r.get("dimensions", {}).get("columns", 0) for r in table_files)
        table_positions = [r.get("position", {}) for r in table_files if r.get("position")]

        # æ„å»ºå…ƒæ•°æ®
        metadata = {
            "processing_date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "source_file": base_filename,
            "source_path": file_path,
            "file_type": "PPTX",
            "processing_method": "advanced_zip_xml_parsing",
            "statistics": {
                "total_slides": total_slides,
                "slides_with_text": len(text_files),
                "slides_with_tables": len([r for r in table_files if r.get("dimensions", {}).get("rows", 0) > 0]),
                "total_tables": total_tables,
                "total_table_rows": total_table_rows,
                "total_table_columns": total_table_columns,
                "slides_with_images": slides_with_images,
                "total_images_extracted": len(image_files_traditional) + len(image_files_zip),
                "images_via_traditional": len(image_files_traditional),
                "images_via_zip_parsing": len(image_files_zip),
                "total_images_in_media": total_images_zip
            },
            "slide_image_mapping": slide_image_mapping,
            "files": {
                "text_files": text_files,
                "table_files": table_files,
                "image_files_traditional": image_files_traditional,
                "image_files_zip": image_files_zip
            },
            "table_analysis": {
                "positions": table_positions,
                "position_stats": {
                    "avg_left": sum(p.get("left", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0,
                    "avg_top": sum(p.get("top", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0,
                    "avg_width": sum(p.get("width", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0,
                    "avg_height": sum(p.get("height", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0
                }
            },
            "processing_info": {
                "extraction_methods": ["python-pptx", "zip_xml_parsing"],
                "image_formats_supported": ["PNG", "WMF", "EMF", "JPEG"],
                "table_extraction_enhanced": True,
                "table_position_tracking": True,
                "clip_descriptions_generated": True,
                "output_format": "JSON/CSV"
            }
        }

        # ä¿å­˜å…ƒæ•°æ®
        metadata_path = f"output/{base_filename}_pptx_metadata.json"
        with open(metadata_path, "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)

        return metadata_path


def process_pptx_file_advanced(preprocessor, file_path, fast_mode=False):
    """
    é«˜çº§PPTXå¤„ç†çš„å…¥å£å‡½æ•°

    Args:
        preprocessor: MultimodalPreprocessorå®ä¾‹
        file_path: PPTXæ–‡ä»¶è·¯å¾„
        fast_mode: å¿«é€Ÿæ¨¡å¼ï¼Œè·³è¿‡è€—æ—¶å¤„ç†
    """
    processor = AdvancedPPTProcessor(preprocessor, fast_mode=fast_mode)
    processor.process_pptx_file_advanced(file_path)


def build_multimodal_knowledge_graph(
        neo4j_uri: str,
        neo4j_user: str,
        neo4j_password: str,
        deepseek_api_key: str,
        input_dir: str = "input",
        output_dir: str = "output",
        document_name: str = "å¤šæ¨¡æ€æ–‡æ¡£",
        fast_mode: bool = False,
        clear_database: bool = False,
        verbose: bool = True
) -> dict:
    """
    æ„å»ºå¤šæ¨¡æ€PPT/PDFçŸ¥è¯†å›¾è°±çš„ä¸»å‡½æ•°

    Args:
        neo4j_uri: Neo4jæ•°æ®åº“URIï¼Œå¦‚ "bolt://localhost:7687"
        neo4j_user: Neo4jç”¨æˆ·å
        neo4j_password: Neo4jå¯†ç 
        deepseek_api_key: DeepSeek APIå¯†é’¥
        input_dir: è¾“å…¥æ–‡ä»¶ç›®å½•ï¼Œé»˜è®¤ä¸º "input"
        output_dir: è¾“å‡ºæ–‡ä»¶ç›®å½•ï¼Œé»˜è®¤ä¸º "output"
        document_name: æ–‡æ¡£åç§°ï¼ˆç”¨äºNeo4jä¸­çš„æ–‡æ¡£èŠ‚ç‚¹ï¼‰ï¼Œé»˜è®¤ä¸º "å¤šæ¨¡æ€æ–‡æ¡£"
        fast_mode: æ˜¯å¦å¯ç”¨å¿«é€Ÿæ¨¡å¼ï¼ˆè·³è¿‡CLIPæè¿°ç”Ÿæˆï¼‰ï¼Œé»˜è®¤ä¸º False
        clear_database: æ˜¯å¦æ¸…ç©ºæ•°æ®åº“ï¼ˆè°¨æ…ä½¿ç”¨ï¼‰ï¼Œé»˜è®¤ä¸º False
        verbose: æ˜¯å¦æ˜¾ç¤ºè¯¦ç»†è¾“å‡ºï¼Œé»˜è®¤ä¸º True

    Returns:
        dict: åŒ…å«å¤„ç†ç»“æœçš„å­—å…¸
    """

    def log(message: str):
        """æ¡ä»¶æ—¥å¿—è¾“å‡º"""
        if verbose:
            print(message)

    # åˆå§‹åŒ–ç»“æœå­—å…¸
    result = {
        'success': False,
        'error': None,
        'statistics': {},
        'files_processed': [],
        'entities_extracted': 0,
        'relationships_extracted': 0,
        'entities_saved': 0,
        'relationships_saved': 0,
        'neo4j_stats': {}
    }

    try:
        if verbose:
            log("=" * 80)
            log("ğŸš€ å¤šæ¨¡æ€PPTçŸ¥è¯†å›¾è°±æ„å»ºç³»ç»Ÿ")
            log("=" * 80)
            log(f"ğŸ“‹ é…ç½®ä¿¡æ¯:")
            log(f"   ğŸ—„ï¸  Neo4j: {neo4j_uri}")
            log(f"   ğŸ¤– DeepSeek API: {'å·²é…ç½®' if deepseek_api_key else 'æœªé…ç½®'}")
            log(f"   ğŸ“ è¾“å…¥ç›®å½•: {input_dir}")
            log(f"   ğŸ“ è¾“å‡ºç›®å½•: {output_dir}")
            log(f"   âš¡ å¿«é€Ÿæ¨¡å¼: {'å¼€å¯' if fast_mode else 'å…³é—­'}")
            log(f"   ğŸ“„ æ–‡æ¡£åç§°: {document_name}")
            log("=" * 80)

        # ==================== ç¬¬ä¸€æ­¥ï¼šåˆå§‹åŒ–å¤šæ¨¡æ€é¢„å¤„ç†å™¨ ====================
        log("\nğŸ”§ ç¬¬ä¸€æ­¥ï¼šåˆå§‹åŒ–å¤šæ¨¡æ€é¢„å¤„ç†å™¨...")
        if verbose:
            log("âš ï¸  æ³¨æ„ï¼šé¦–æ¬¡è¿è¡Œå¯èƒ½éœ€è¦ä¸‹è½½æ¨¡å‹ï¼Œè¯·è€å¿ƒç­‰å¾…...")

        processor = MultimodalPreprocessor()
        log("âœ… å¤šæ¨¡æ€é¢„å¤„ç†å™¨åˆå§‹åŒ–å®Œæˆ")

        # ==================== ç¬¬äºŒæ­¥ï¼šæ£€æŸ¥å¹¶å¤„ç†è¾“å…¥æ–‡ä»¶ ====================
        log(f"\nğŸ“‚ ç¬¬äºŒæ­¥ï¼šæ£€æŸ¥è¾“å…¥ç›®å½• {input_dir}...")

        # ç¡®ä¿è¾“å…¥ç›®å½•å­˜åœ¨
        if not os.path.exists(input_dir):
            os.makedirs(input_dir, exist_ok=True)
            error_msg = f"è¾“å…¥ç›®å½•ä¸å­˜åœ¨ï¼Œå·²åˆ›å»º: {input_dir}ï¼Œè¯·å°†PPT/PDFæ–‡ä»¶æ”¾å…¥è¯¥ç›®å½•"
            result['error'] = error_msg
            log(f"âŒ {error_msg}")
            return result

        # æŸ¥æ‰¾æ”¯æŒçš„æ–‡ä»¶
        input_files = [f for f in os.listdir(input_dir)
                       if f.lower().endswith(('.pdf', '.pptx')) and not f.startswith('~$')]

        if not input_files:
            error_msg = f"åœ¨ {input_dir} ç›®å½•ä¸­æœªæ‰¾åˆ°PPT/PDFæ–‡ä»¶"
            result['error'] = error_msg
            log(f"âŒ {error_msg}")
            return result

        ppt_count = sum(1 for f in input_files if f.lower().endswith('.pptx'))
        pdf_count = sum(1 for f in input_files if f.lower().endswith('.pdf'))
        log(f"âœ… æ‰¾åˆ°æ–‡ä»¶: {ppt_count}ä¸ªPPT, {pdf_count}ä¸ªPDF")

        result['files_processed'] = input_files
        result['statistics']['ppt_count'] = ppt_count
        result['statistics']['pdf_count'] = pdf_count

        # ==================== ç¬¬ä¸‰æ­¥ï¼šå¤„ç†æ–‡ä»¶ï¼ˆå¤šæ¨¡æ€é¢„å¤„ç†ï¼‰====================
        log(f"\nğŸ”„ ç¬¬ä¸‰æ­¥ï¼šå¤šæ¨¡æ€æ•°æ®é¢„å¤„ç†...")

        for idx, input_file in enumerate(input_files, 1):
            input_path = os.path.join(input_dir, input_file)
            log(f"\nğŸ“„ [{idx}/{len(input_files)}] å¤„ç†æ–‡ä»¶: {input_file}")

            if input_file.lower().endswith('.pdf'):
                log("   ğŸ“š ä½¿ç”¨PDFå¤„ç†å™¨...")
                processor.process_pdf(input_path)

            elif input_file.lower().endswith('.pptx'):
                log("   ğŸ“Š ä½¿ç”¨é«˜çº§PPTå¤„ç†å™¨...")
                # åˆ›å»ºé«˜çº§PPTå¤„ç†å™¨
                ppt_processor = AdvancedPPTProcessor(processor, fast_mode=fast_mode)
                ppt_processor.process_pptx_file_advanced(input_path)

            log(f"   âœ… [{idx}/{len(input_files)}] å®Œæˆ: {input_file}")

        log(f"\nğŸ‰ å¤šæ¨¡æ€é¢„å¤„ç†å®Œæˆï¼å…±å¤„ç† {len(input_files)} ä¸ªæ–‡ä»¶")
        log(f"ğŸ“ é¢„å¤„ç†ç»“æœä¿å­˜åœ¨: {output_dir}/")

        # ==================== ç¬¬å››æ­¥ï¼šå®ä½“å…³ç³»æŠ½å– ====================
        log(f"\nğŸ” ç¬¬å››æ­¥ï¼šå®ä½“å…³ç³»æŠ½å–...")

        if not deepseek_api_key:
            error_msg = "DeepSeek API Keyæœªé…ç½®ï¼Œæ— æ³•è¿›è¡Œå®ä½“æŠ½å–"
            result['error'] = error_msg
            log(f"âŒ {error_msg}")
            return result

        log("   ğŸ¤– åˆå§‹åŒ–å®ä½“æŠ½å–å™¨...")
        extractor = EntityExtractor(deepseek_api_key)

        log("   ğŸ“¥ åŠ è½½å¤šæ¨¡æ€æ•°æ®...")
        multimodal_data = extractor.load_multimodal_data(output_dir)

        log(f"   ğŸ“Š æ•°æ®ç»Ÿè®¡: {len(multimodal_data.get('slides', []))}ä¸ªå¹»ç¯ç‰‡, {len(multimodal_data.get('images', []))}å¼ å›¾ç‰‡")

        log("   ğŸ”¬ å¼€å§‹å®ä½“å…³ç³»æŠ½å–...")
        extracted_data = extractor.extract_entities_from_multimodal(multimodal_data)

        result['entities_extracted'] = len(extracted_data.entities)
        result['relationships_extracted'] = len(extracted_data.relationships)

        log(f"âœ… å®ä½“å…³ç³»æŠ½å–å®Œæˆ:")
        log(f"   ğŸ·ï¸  å®ä½“æ•°é‡: {len(extracted_data.entities)}")
        log(f"   ğŸ”— å…³ç³»æ•°é‡: {len(extracted_data.relationships)}")
        log(f"   ğŸ“‹ å±æ€§æ•°é‡: {len(extracted_data.attributes)}")

        # ==================== ç¬¬äº”æ­¥ï¼šä¿å­˜åˆ°Neo4jçŸ¥è¯†å›¾è°± ====================
        log(f"\nğŸ’¾ ç¬¬äº”æ­¥ï¼šä¿å­˜åˆ°Neo4jçŸ¥è¯†å›¾è°±...")

        log("   ğŸ”Œ è¿æ¥Neo4jæ•°æ®åº“...")
        kg = Neo4jKnowledgeGraph(neo4j_uri, neo4j_user, neo4j_password)

        # å¯é€‰ï¼šæ¸…ç©ºæ•°æ®åº“
        if clear_database:
            log("   ğŸ—‘ï¸  æ¸…ç©ºæ•°æ®åº“...")
            kg.clear_database()

        log("   ğŸ’¾ ä¿å­˜å®ä½“å…³ç³»æ•°æ®...")
        save_result = kg.save_extracted_data(extracted_data, document_name)

        if save_result['success']:
            result['entities_saved'] = save_result['entities_saved']
            result['relationships_saved'] = save_result['relationships_saved']

            log(f"âœ… æ•°æ®ä¿å­˜æˆåŠŸ:")
            log(f"   ğŸ“„ æ–‡æ¡£: {save_result['document']}")
            log(f"   ğŸ·ï¸  å®ä½“: {save_result['entities_saved']}/{len(extracted_data.entities)}")
            log(f"   ğŸ”— å…³ç³»: {save_result['relationships_saved']}/{len(extracted_data.relationships)}")
        else:
            error_msg = f"æ•°æ®ä¿å­˜å¤±è´¥: {save_result.get('error', 'æœªçŸ¥é”™è¯¯')}"
            result['error'] = error_msg
            log(f"âŒ {error_msg}")
            kg.close()
            return result

        # ==================== ç¬¬å…­æ­¥ï¼šæ˜¾ç¤ºçŸ¥è¯†å›¾è°±ç»Ÿè®¡ä¿¡æ¯ ====================
        log(f"\nğŸ“Š ç¬¬å…­æ­¥ï¼šçŸ¥è¯†å›¾è°±ç»Ÿè®¡ä¿¡æ¯...")

        stats = kg.get_statistics()
        result['neo4j_stats'] = stats

        log(f"   ğŸ“„ æ–‡æ¡£æ€»æ•°: {stats['total_documents']}")
        log(f"   ğŸ·ï¸  èŠ‚ç‚¹æ€»æ•°: {stats['total_nodes']}")
        log(f"   ğŸ”— å…³ç³»æ€»æ•°: {stats['total_relationships']}")

        if verbose:
            log(f"\nğŸ·ï¸  å®ä½“ç±»å‹åˆ†å¸ƒ:")
            for entity_type in stats['entity_types'][:10]:  # æ˜¾ç¤ºå‰10ç§ç±»å‹
                log(f"   - {entity_type['entity_type']}: {entity_type['count']}ä¸ª")

            # ==================== ç¬¬ä¸ƒæ­¥ï¼šæ˜¾ç¤ºæŸ¥è¯¢ç¤ºä¾‹ ====================
            log(f"\nğŸ” ç¬¬ä¸ƒæ­¥ï¼šçŸ¥è¯†å›¾è°±æŸ¥è¯¢ç¤ºä¾‹...")

            # æŸ¥è¯¢å®ä½“ç¤ºä¾‹
            log(f"\nğŸ·ï¸  å®ä½“ç¤ºä¾‹ (å‰5ä¸ª):")
            entities = kg.query_entities(5)
            for i, entity in enumerate(entities, 1):
                log(f"   {i}. {entity['name']} ({entity['type']}) - {entity.get('description', '')[:50]}...")

            # æŸ¥è¯¢å…³ç³»ç¤ºä¾‹
            log(f"\nğŸ”— å…³ç³»ç¤ºä¾‹ (å‰5ä¸ª):")
            relationships = kg.query_relationships(5)
            for i, rel in enumerate(relationships, 1):
                log(f"   {i}. {rel['source']} --{rel['relation']}--> {rel['target']}")

            # æœç´¢ç‰¹å®šå®ä½“ç¤ºä¾‹
            if entities:
                sample_entity = entities[0]['name']
                log(f"\nğŸ” å®ä½“é‚»å±…æŸ¥è¯¢ç¤ºä¾‹ (ä»¥ '{sample_entity}' ä¸ºä¾‹):")
                neighbors = kg.get_entity_neighbors(sample_entity, depth=1)

                log(f"   é‚»å±…å®ä½“ (å‰3ä¸ª):")
                for i, neighbor in enumerate(neighbors['neighbors'][:3], 1):
                    log(f"     {i}. {neighbor['name']} ({neighbor['type']})")

                log(f"   ç›¸å…³å…³ç³» (å‰3ä¸ª):")
                for i, rel in enumerate(neighbors['relationships'][:3], 1):
                    direction = "â†’" if rel['direction'] == 'outgoing' else "â†"
                    log(f"     {i}. {sample_entity} {direction}[{rel['relation']}] {rel['neighbor']}")

        # å…³é—­æ•°æ®åº“è¿æ¥
        kg.close()

        # è®¾ç½®æˆåŠŸæ ‡å¿—
        result['success'] = True

        # ==================== å®Œæˆæ€»ç»“ ====================
        if verbose:
            log(f"\n" + "=" * 80)
            log("ğŸ‰ å¤šæ¨¡æ€PPTçŸ¥è¯†å›¾è°±æ„å»ºå®Œæˆï¼")
            log("=" * 80)
            log(f"ğŸ“‹ å¤„ç†æ€»ç»“:")
            log(f"   ğŸ“ è¾“å…¥æ–‡ä»¶: {len(input_files)}ä¸ª ({ppt_count}ä¸ªPPT, {pdf_count}ä¸ªPDF)")
            log(f"   ğŸ”„ å¤šæ¨¡æ€å¤„ç†: âœ… å®Œæˆ")
            log(f"   ğŸ” å®ä½“æŠ½å–: âœ… {len(extracted_data.entities)}ä¸ªå®ä½“, {len(extracted_data.relationships)}ä¸ªå…³ç³»")
            log(f"   ğŸ’¾ Neo4jå­˜å‚¨: âœ… {save_result['entities_saved']}ä¸ªå®ä½“, {save_result['relationships_saved']}ä¸ªå…³ç³»")
            log(f"   ğŸ—„ï¸  æ•°æ®åº“ç»Ÿè®¡: {stats['total_nodes']}ä¸ªèŠ‚ç‚¹, {stats['total_relationships']}ä¸ªå…³ç³»")

            log(f"\nğŸ“ è¾“å‡ºæ–‡ä»¶ä½ç½®:")
            log(f"   - å¤šæ¨¡æ€æ•°æ®: {output_dir}/")
            log(f"   - æ–‡æœ¬æ•°æ®: {output_dir}/text/")
            log(f"   - å›¾ç‰‡æ•°æ®: {output_dir}/images/")
            log(f"   - è¡¨æ ¼æ•°æ®: {output_dir}/tables/")
            log(f"   - å…¬å¼æ•°æ®: {output_dir}/formulas/")
            log(f"   - ä»£ç æ•°æ®: {output_dir}/code/")

            log(f"\nğŸ” Neo4jæŸ¥è¯¢å»ºè®®:")
            log(f"   - æŸ¥çœ‹æ‰€æœ‰å®ä½“: MATCH (e:Entity) RETURN e LIMIT 25")
            log(f"   - æŸ¥çœ‹æ‰€æœ‰å…³ç³»: MATCH (s)-[r]->(t) RETURN s.name, type(r), t.name LIMIT 25")
            log(f"   - æŸ¥çœ‹æ–‡æ¡£: MATCH (d:Document) RETURN d")
            log(f"   - æœç´¢Arduinoç›¸å…³: MATCH (e:Entity) WHERE toLower(e.name) CONTAINS 'arduino' RETURN e")
            log("=" * 80)

        return result

    except KeyboardInterrupt:
        error_msg = "ç”¨æˆ·ä¸­æ–­æ“ä½œ"
        result['error'] = error_msg
        log(f"\nâš ï¸  {error_msg}")
        return result

    except Exception as e:
        error_msg = f"ç³»ç»Ÿé”™è¯¯: {str(e)}"
        result['error'] = error_msg
        log(f"\nâŒ {error_msg}")

        if verbose:
            import traceback
            traceback.print_exc()
            log(f"\nğŸ”§ æ•…éšœæ’é™¤å»ºè®®:")
            log(f"   1. æ£€æŸ¥Neo4jæ•°æ®åº“è¿æ¥")
            log(f"   2. æ£€æŸ¥DeepSeek API Keyæ˜¯å¦æœ‰æ•ˆ")
            log(f"   3. ç¡®ä¿inputç›®å½•ä¸­æœ‰PPT/PDFæ–‡ä»¶")
            log(f"   4. æ£€æŸ¥ç½‘ç»œè¿æ¥ï¼ˆæ¨¡å‹ä¸‹è½½éœ€è¦ï¼‰")
            log(f"   5. æ£€æŸ¥ç£ç›˜ç©ºé—´æ˜¯å¦å……è¶³")

        return result

from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent.parent  # è·å–æ ¹ç›®å½•
input_file_path = BASE_DIR / "input"
output_file_path = BASE_DIR / "models" / "output"

# ä¿ç•™åŸæ¥çš„ä¸»å‡½æ•°ä½œä¸ºç¤ºä¾‹
if __name__ == "__main__":

    result = build_multimodal_knowledge_graph(
        neo4j_uri="bolt://101.132.130.25:7687",
        neo4j_user="neo4j",
        neo4j_password="wangshuxvan@1",
        deepseek_api_key="sk-c28ec338b39e4552b9e6bded47466442",
        input_dir=input_file_path,
        output_dir=output_file_path,
        document_name="Arduinoè¯¾ç¨‹PPT",
        fast_mode=False,
        clear_database=True,
        verbose=True
    )
